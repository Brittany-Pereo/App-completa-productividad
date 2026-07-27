# -*- coding: utf-8 -*-
"""Kit compartido de formas nativas de PowerPoint (editables) y las gráficas
que se construyen con ellas.

Antes esto vivía duplicado casi igual en dos proyectos (planes_de_justicia y
crear_pptx_clues) — cada uno con su propia copia de _forma_rect, _forma_texto,
etc. Como los dos generan el mismo tipo de reporte (mismo diseño, mismos
colores), lo movimos aquí una sola vez.

Nota sobre nombres: las funciones de "bajo nivel" (forma_rect, forma_texto...)
ya NO llevan guion bajo al inicio como en las versiones originales. Un guion
bajo al inicio (`_forma_rect`) es una convención de Python para decir "esto es
un detalle interno del archivo, no lo uses desde afuera". Como ahora este
archivo SÍ está pensado para que otros módulos lo importen, se vuelven
funciones "públicas".
"""

from __future__ import annotations

import math
from datetime import date

import numpy as np
import pandas as pd
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

# ---------------------------------------------------------------------------
# Paleta de colores compartida
# ---------------------------------------------------------------------------
COL_MUTED = "#6B7280"
COL_TEXTO = "#111827"
COL_BORDE = "#D1D5DB"
COL_VERDE = "#1E5B4F"
COL_GUINDA = "#611232"
COL_DORADO = "#A57F2C"

MESES_ES = [
    "enero", "febrero", "marzo", "abril", "mayo", "junio",
    "julio", "agosto", "septiembre", "octubre", "noviembre", "diciembre",
]
MESES_EN_ABBR = ["Jan", "Feb", "Mar", "Apr", "May", "Jun",
                  "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
MESES_ES_ABBR = ["ene", "feb", "mar", "abr", "may", "jun",
                  "jul", "ago", "sep", "oct", "nov", "dic"]


def rgb(hex_color: str) -> RGBColor:
    """Convierte '#RRGGBB' a RGBColor de python-pptx."""
    return RGBColor.from_string(hex_color.lstrip("#").upper())


def fmt_num(x) -> str:
    """1234567 -> '1,234,567' (equivalente a scales::comma en R)."""
    try:
        return f"{int(round(float(x))):,}"
    except (TypeError, ValueError):
        return "0"


def floor_month(d) -> date:
    d = pd.to_datetime(d).date() if not isinstance(d, date) else d
    return date(d.year, d.month, 1)


def ceiling_month(d) -> date:
    d = pd.to_datetime(d).date() if not isinstance(d, date) else d
    if d.month == 12:
        return date(d.year + 1, 1, 1)
    return date(d.year, d.month + 1, 1)


def mes_titulo(ts) -> str:
    ts = pd.to_datetime(ts)
    return f"{MESES_ES[ts.month - 1].capitalize()} {ts.year}"


def mes_abbr_titulo(ts) -> str:
    ts = pd.to_datetime(ts)
    return f"{MESES_ES[ts.month - 1][:3].capitalize()}-{ts.year}"


# ---------------------------------------------------------------------------
# Formas nativas de bajo nivel
# ---------------------------------------------------------------------------
def forma_rect(slide, left, top, width, height, color_hex, transparencia_pct=0):
    """Rectángulo nativo relleno. `transparencia_pct`: 0 (opaco) a 100 (invisible)."""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, int(left), int(top), int(max(width, 1)), int(max(height, 1)))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(color_hex)
    if transparencia_pct:
        srgb = shp.fill._xPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
        alpha = srgb.makeelement(qn("a:alpha"), {"val": str(int((100 - transparencia_pct) * 1000))})
        srgb.append(alpha)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def forma_texto(slide, left, top, width, height, texto, size=9, color_hex=COL_TEXTO,
                bold=False, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                rotation=0, wrap=True, font="Calibri"):
    tb = slide.shapes.add_textbox(int(left), int(top), int(max(width, 1)), int(max(height, 1)))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, linea in enumerate(str(texto).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = linea
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = rgb(color_hex)
        r.font.name = font
    if rotation:
        tb.rotation = rotation
    return tb


def forma_ovalo(slide, cx, cy, radio, color_hex):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, int(cx - radio), int(cy - radio), int(radio * 2), int(radio * 2))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(color_hex)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def forma_linea_quebrada(slide, puntos, color_hex, width_pt=1.5):
    """Línea poligonal nativa y editable (freeform) que conecta `puntos`
    (lista de (x_emu, y_emu)) — se puede editar punto por punto en PowerPoint."""
    if len(puntos) < 2:
        return None
    fb = slide.shapes.build_freeform(int(puntos[0][0]), int(puntos[0][1]), scale=1.0)
    fb.add_line_segments([(int(x), int(y)) for x, y in puntos[1:]], close=False)
    shp = fb.convert_to_shape()
    shp.fill.background()
    shp.line.color.rgb = rgb(color_hex)
    shp.line.width = Pt(width_pt)
    shp.shadow.inherit = False
    return shp


def forma_flecha_vertical(slide, x, y_top, y_bottom, color_hex, width_pt=1.2):
    """Conector vertical con punta de flecha arriba (en `y_top`)."""
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, int(x), int(y_top), int(x), int(y_bottom))
    conn.line.color.rgb = rgb(color_hex)
    conn.line.width = Pt(width_pt)
    ln = conn.line._get_or_add_ln()
    head = ln.makeelement(qn("a:headEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(head)
    return conn


def forma_linea_punteada(slide, x1, y1, x2, y2, color_hex, width_pt=0.75):
    """Conector recto punteado (equivalente a geom_segment(linetype='dotted')),
    usado como línea guía entre un punto destacado y su etiqueta."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, int(x1), int(y1), int(x2), int(y2))
    conn.line.color.rgb = rgb(color_hex)
    conn.line.width = Pt(width_pt)
    conn.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT
    conn.shadow.inherit = False
    return conn


def escala_bonita(valor_max, n_pasos_objetivo=5):
    """Devuelve (paso, valor_redondeado) con incrementos 'bonitos' (1/2/2.5/5/10 x 10^n),
    igual que el locator automático de matplotlib, para que el eje Y no muestre
    números arbitrarios como 9,420,922."""
    if valor_max <= 0:
        return 1, n_pasos_objetivo
    bruto = valor_max / n_pasos_objetivo
    exponente = math.floor(math.log10(bruto))
    base = 10 ** exponente
    paso = base * 10
    for m in (1, 2, 2.5, 5, 10):
        if bruto <= m * base:
            paso = m * base
            break
    return paso, math.ceil(valor_max / paso) * paso


def dibujar_eje_y(slide, plot_l, plot_w, baseline, plot_h, ymax_eje, valores_marca,
                  ancho_etiqueta=None, gridlines=False, size=8, ymin_eje=0):
    """Dibuja las marcas del eje Y (texto, en `valores_marca`) y opcionalmente
    líneas de rejilla, mapeadas proporcionalmente contra el rango
    [`ymin_eje`, `ymax_eje`]."""
    ancho_etiqueta = ancho_etiqueta or Emu(900000)
    alto_etiqueta = Pt(size + 6)
    rango_eje = ymax_eje - ymin_eje
    for valor in valores_marca:
        frac = ((valor - ymin_eje) / rango_eje) if rango_eje else 0
        y = baseline - frac * plot_h
        forma_texto(
            slide, plot_l - ancho_etiqueta - Pt(4), y - alto_etiqueta / 2, ancho_etiqueta, alto_etiqueta,
            fmt_num(valor), size=size, color_hex=COL_MUTED, align=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        if gridlines and valor > ymin_eje:
            forma_rect(slide, plot_l, y, plot_w, Pt(0.6), "#E5E7EB")


# ---------------------------------------------------------------------------
# Gráficas de barras (histórico y avance-vs-meta) — formas nativas
# ---------------------------------------------------------------------------
def dibujar_grafica_barras(slide, box, categorias, totales, avances, titulo,
                           colores_total, colores_avance,
                           etiquetas_total=None, etiquetas_avance=None,
                           title_size=15):
    """Barras de 'total' (claro) con 'avance' superpuesto (oscuro), con
    etiquetas de valor."""
    L, T, W, H = box
    n = len(categorias)
    if n == 0:
        return

    if etiquetas_total is None:
        etiquetas_total = [fmt_num(t) for t in totales]
    if etiquetas_avance is None:
        etiquetas_avance = [fmt_num(a) for a in avances]

    alto_titulo = int(H * 0.13)
    alto_categoria = Pt(16)
    margen_izq = int(W * 0.12)
    espacio_etiqueta_sup = int(H * 0.16)

    forma_texto(slide, L, T, W, alto_titulo, titulo, size=title_size,
               color_hex=COL_MUTED, bold=True, align=PP_ALIGN.CENTER)

    plot_l = L + margen_izq
    plot_t = T + alto_titulo + espacio_etiqueta_sup
    plot_w = W - margen_izq
    plot_h = H - alto_titulo - espacio_etiqueta_sup - alto_categoria
    baseline = plot_t + plot_h

    ymax_datos = max(max(totales, default=0), max(avances, default=0), 1)
    paso, ymax_redondeado = escala_bonita(ymax_datos, 5)
    ymax_eje = max(ymax_redondeado, ymax_datos) * 1.16

    valores_marca = [i * paso for i in range(int(ymax_eje // paso) + 1)]
    dibujar_eje_y(slide, plot_l, plot_w, baseline, plot_h, ymax_eje, valores_marca,
                 ancho_etiqueta=margen_izq - Pt(4))

    slot_w = plot_w / n
    bar_w = slot_w * 0.60

    for i, cat in enumerate(categorias):
        slot_l = plot_l + i * slot_w
        bar_l = slot_l + (slot_w - bar_w) / 2

        h_total = (totales[i] / ymax_eje) * plot_h if ymax_eje else 0
        h_avance = (avances[i] / ymax_eje) * plot_h if ymax_eje else 0

        if h_total > 0:
            forma_rect(slide, bar_l, baseline - h_total, bar_w, h_total, colores_total[i])
        if h_avance > 0:
            forma_rect(slide, bar_l, baseline - h_avance, bar_w, h_avance, colores_avance[i])

        n_lineas_tot = etiquetas_total[i].count("\n") + 1
        alto_et = Pt(13) * n_lineas_tot
        forma_texto(slide, slot_l, baseline - h_total - alto_et - Pt(3), slot_w, alto_et,
                   etiquetas_total[i], size=9.5, color_hex="#000000", bold=True,
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)

        if h_avance > Pt(24):
            n_lineas_av = etiquetas_avance[i].count("\n") + 1
            alto_ea = Pt(12) * n_lineas_av
            forma_texto(slide, slot_l, baseline - h_avance + Pt(3), slot_w, alto_ea,
                       etiquetas_avance[i], size=8.5, color_hex="#FFFFFF", bold=True,
                       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

        forma_texto(slide, slot_l, baseline + Pt(3), slot_w, alto_categoria,
                   str(cat), size=11, color_hex=COL_MUTED, bold=True, align=PP_ALIGN.CENTER)


def dibujar_grafica_planeacion_historica(slide, box, df, col_total, col_avance, titulo,
                                         beige="#D9D2BE", verde="#2F6F63"):
    """Barras 2020-2025: total (beige) con avance (verde) superpuesto."""
    anios = list(range(2020, 2026))
    d = df.copy()
    d["anio_num"] = pd.to_numeric(d["anio"], errors="coerce").astype("Int64")
    d = d[d["anio_num"].isin(anios)]
    d = d.set_index("anio_num").reindex(anios)
    totales = d[col_total].fillna(0).astype(float).tolist()
    avances = d[col_avance].fillna(0).astype(float).tolist()

    dibujar_grafica_barras(
        slide, box, categorias=[str(a) for a in anios],
        totales=totales, avances=avances, titulo=titulo,
        colores_total=[beige] * len(anios), colores_avance=[verde] * len(anios),
    )


def dibujar_grafica_planeacion_2024_2026(slide, box, df, col_total, col_avance, titulo,
                                         beige="#D9D2BE", verde="#2F6F63",
                                         beige_2026="#A99F86", verde_2026="#1E5B4F"):
    """Barras 2024-2026 con etiquetas especiales de 'Meta 2026' / 'Avance'."""
    anios = [2024, 2025, 2026]
    d = df.copy()
    d["anio_num"] = pd.to_numeric(d["anio"], errors="coerce").astype("Int64")
    d = d[d["anio_num"].isin(anios)]
    d = d.set_index("anio_num").reindex(anios)
    totales = d[col_total].fillna(0).astype(float).tolist()
    avances = d[col_avance].fillna(0).astype(float).tolist()

    etiquetas_total, etiquetas_avance = [], []
    for a, tot, av in zip(anios, totales, avances):
        if a == 2026:
            etiquetas_total.append(f"Meta 2026\n{fmt_num(tot)}")
            pct = (av / tot) if tot > 0 else np.nan
            pct_txt = "s/d" if pd.isna(pct) else f"{int(round(pct * 100))}%"
            etiquetas_avance.append(f"Avance\n{fmt_num(av)}\n({pct_txt})")
        else:
            etiquetas_total.append(fmt_num(tot))
            etiquetas_avance.append(fmt_num(av))

    colores_total = [beige_2026 if a == 2026 else beige for a in anios]
    colores_avance = [verde_2026 if a == 2026 else verde for a in anios]

    dibujar_grafica_barras(
        slide, box, categorias=[str(a) for a in anios],
        totales=totales, avances=avances, titulo=titulo,
        colores_total=colores_total, colores_avance=colores_avance,
        etiquetas_total=etiquetas_total, etiquetas_avance=etiquetas_avance,
    )


# ---------------------------------------------------------------------------
# Gráfica de serie temporal (consultas/procedimientos por mes) — formas nativas
# ---------------------------------------------------------------------------
def dibujar_grafica_consultas_periodos(slide, box, df, fecha_inicio="2022-08-01", fecha_fin=None,
                                       titulo=None, title_size=17,
                                       color_linea="#6B6B6B", verde_punto="#1F5B50",
                                       fill_2223="#EFEFEF", fill_2024="#E9DDCC",
                                       fill_2025="#F4F0EA", fill_2026="#E9DDCC",
                                       fill_valuebox="#B99C6D"):
    """Serie temporal mensual con bandas por periodo y anotaciones, dibujada
    con formas nativas de PowerPoint (línea freeform editable, rectángulos,
    óvalos y textos) en vez de una imagen.

    `titulo`: si se pasa, se dibuja dentro del área del gráfico (arriba). Si
    es None, no se reserva espacio para título (úsalo cuando el título ya
    va en un placeholder de texto aparte, como el "fecha" de la diapositiva).
    """
    L, T, W, H = box
    hoy_mes = floor_month(date.today())
    if fecha_fin is None:
        fecha_fin = hoy_mes
    else:
        fecha_fin = pd.to_datetime(fecha_fin).date()
    fecha_inicio = pd.to_datetime(fecha_inicio).date()

    alto_titulo = Pt(24) if titulo else 0
    if titulo:
        forma_texto(slide, L, T, W, alto_titulo, titulo, size=title_size, color_hex=COL_TEXTO,
                   bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    d = df.copy()
    d["fecha"] = pd.to_datetime(d["fecha"]).dt.tz_localize(None)
    d = d[(d["fecha"].dt.date >= fecha_inicio) &
          (d["fecha"].dt.date <= fecha_fin) &
          (d["fecha"].dt.date < hoy_mes)].sort_values("fecha")

    if len(d) == 0:
        forma_texto(slide, L, T + alto_titulo, W, H - alto_titulo,
                   "Sin datos suficientes para este período", size=12, color_hex=COL_MUTED)
        return

    ymax = float(d["consultas_totales"].max())
    ymin = float(d["consultas_totales"].min())
    ymin_eje = max(ymin * 0.90, 0)
    ymax_eje = max(ymax * 1.10, ymin_eje + 1)

    fecha_fin_banda = ceiling_month(fecha_fin)
    bandas = [
        (fecha_inicio, date(2024, 1, 1), fill_2223, "2022–2023\nAños de transición"),
        (date(2024, 1, 1), date(2025, 1, 1), fill_2024, "2024\nPrimer año de operación"),
        (date(2025, 1, 1), date(2026, 1, 1), fill_2025, "2025\nSegundo año de operación"),
        (date(2026, 1, 1), fecha_fin_banda, fill_2026, "2026\nTercer año de operación"),
    ]

    dias_pad = max((fecha_fin_banda - fecha_inicio).days * 0.035, 10)
    x0_ord = fecha_inicio.toordinal() - dias_pad
    x1_ord = fecha_fin_banda.toordinal() + dias_pad

    margen_izq = int(W * 0.08)
    margen_der = int(W * 0.02)
    margen_sup = alto_titulo + int(H * (0.03 if titulo else 0.03))
    margen_inf = int(H * 0.14)

    plot_l = L + margen_izq
    plot_t = T + margen_sup
    plot_w = W - margen_izq - margen_der
    plot_h = H - margen_sup - margen_inf
    baseline = plot_t + plot_h

    def _ord(valor_fecha):
        if hasattr(valor_fecha, "toordinal"):
            return valor_fecha.toordinal()
        return pd.Timestamp(valor_fecha).toordinal()

    def xmap(valor_fecha):
        frac = (_ord(valor_fecha) - x0_ord) / (x1_ord - x0_ord)
        return plot_l + frac * plot_w

    def ymap(valor):
        frac = ((valor - ymin_eje) / (ymax_eje - ymin_eje)) if (ymax_eje - ymin_eje) else 0
        return baseline - frac * plot_h

    # Fondo blanco del área de la gráfica
    forma_rect(slide, L, T + alto_titulo, W, H - alto_titulo, "#FFFFFF")

    # Bandas de periodo
    for xmin_b, xmax_b, fill, _lab in bandas:
        x_l = xmap(xmin_b)
        x_r = xmap(xmax_b)
        forma_rect(slide, x_l, plot_t, x_r - x_l, plot_h, fill)

    # Rejilla horizontal + etiquetas del eje Y
    paso_eje, _ = escala_bonita(ymax_eje - ymin_eje, 5)
    primer_marca = math.ceil(ymin_eje / paso_eje) * paso_eje
    valores_marca_y = []
    v = primer_marca
    while v <= ymax_eje:
        valores_marca_y.append(v)
        v += paso_eje
    dibujar_eje_y(slide, plot_l, plot_w, baseline, plot_h, ymax_eje, valores_marca_y,
                 ancho_etiqueta=margen_izq - Pt(4), gridlines=True, size=10, ymin_eje=ymin_eje)

    # Serie: línea nativa (freeform, editable) + marcador en cada punto
    puntos_linea = [(xmap(f.date()), ymap(v)) for f, v in zip(d["fecha"], d["consultas_totales"])]
    forma_linea_quebrada(slide, puntos_linea, color_linea, width_pt=1.5)
    for x, y in puntos_linea:
        forma_ovalo(slide, x, y, Pt(1.8), color_linea)

    # Puntos destacados (mismo mes que el corte, años anteriores a 2026)
    mes_destacado = fecha_fin.month
    puntos_destacados = d[(d["fecha"].dt.month == mes_destacado) &
                          (d["fecha"].dt.year < 2026)].drop_duplicates("fecha")
    for _, row in puntos_destacados.iterrows():
        anio = row["fecha"].year
        valor = row["consultas_totales"]

        # Ajuste fino por año para que las etiquetas no se encimen.
        if anio == 2023:
            nudge_y = -ymax * 0.04
        elif anio == 2024:
            nudge_y = -ymax * 0.05
        elif anio == 2025:
            nudge_y = -ymax * 0.09
        else:
            nudge_y = -ymax * 0.05
        nudge_x_dias = -18 if anio == 2025 else 0

        x_punto, y_punto = xmap(row["fecha"].date()), ymap(valor)
        forma_ovalo(slide, x_punto, y_punto, Pt(5.5), verde_punto)

        fecha_etiqueta = row["fecha"] + pd.Timedelta(days=nudge_x_dias)
        x_etiqueta = xmap(fecha_etiqueta.date())
        y_etiqueta = ymap(valor + nudge_y)

        # Línea punteada del punto hacia la etiqueta
        y_seg_ini = ymap(valor - ymax * 0.02)
        y_seg_fin = ymap(valor + nudge_y + ymax * 0.015)
        forma_linea_punteada(slide, x_punto, y_seg_ini, x_etiqueta, y_seg_fin, "#111827", width_pt=0.75)

        etiqueta = f"{fmt_num(valor)}\n{MESES_ES[row['fecha'].month - 1]} {anio}"
        forma_texto(slide, x_etiqueta - Pt(52), y_etiqueta - Pt(17), Pt(104), Pt(34), etiqueta,
                   size=10.5, color_hex="#000000", bold=True,
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Fila de encabezado (títulos de banda + anotación del decreto) anclada al
    # techo fijo del área de gráfica, para que nunca choque con el recuadro
    # del último valor (su altura depende de dónde caiga el último dato).
    header_y = plot_t + Pt(4)

    # Texto superior de bandas
    for xmin_b, xmax_b, _fill, lab in bandas:
        centro = xmap(xmin_b) + (xmap(xmax_b) - xmap(xmin_b)) / 2
        forma_texto(slide, centro - Pt(80), header_y, Pt(160), Pt(32), lab,
                   size=11, color_hex="#000000", bold=True,
                   align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    # Flecha "Decreto de creación"
    x_decreto = xmap(date(2022, 8, 15))
    forma_flecha_vertical(slide, x_decreto, header_y, baseline, verde_punto)
    forma_texto(slide, x_decreto + Pt(6), header_y, Pt(130), Pt(30),
               "Decreto de creación\ndel IMSS Bienestar", size=9.5, color_hex="#000000",
               bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

    # Recuadro (valuebox) del último dato
    fecha_ultimo = d["fecha"].max()
    valor_ultimo_ser = d.loc[d["fecha"] == fecha_ultimo, "consultas_totales"]
    valor_ultimo = valor_ultimo_ser.iloc[0] if len(valor_ultimo_ser) else 0
    if pd.isna(valor_ultimo):
        valor_ultimo = 0
    x_vb = xmap(fecha_ultimo.date())
    y_vb = ymap(valor_ultimo)
    vb_w, vb_h = Pt(88), Pt(36)
    # No dejar que el recuadro invada la fila de encabezado (banda + decreto).
    vb_top = max(y_vb - vb_h - Pt(10), header_y + Pt(38))
    vb = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, int(x_vb - vb_w / 2), int(vb_top),
        int(vb_w), int(vb_h))
    vb.fill.solid()
    vb.fill.fore_color.rgb = rgb(fill_valuebox)
    vb.line.fill.background()
    vb.shadow.inherit = False
    tf = vb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lineas_valuebox = [fmt_num(valor_ultimo), f"{MESES_ES[fecha_ultimo.month - 1]} {fecha_ultimo.year}"]
    for i, linea in enumerate(lineas_valuebox):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = linea
        r.font.size = Pt(11.5)
        r.font.bold = True
        r.font.color.rgb = rgb("#FFFFFF")
        r.font.name = "Calibri"

    # Etiquetas del eje X (una marca por mes)
    eje_inicio = pd.Timestamp(fecha_inicio).replace(day=1)
    eje_fin = pd.Timestamp(fecha_fin_banda)
    for tick in pd.date_range(eje_inicio, eje_fin, freq="MS"):
        x = xmap(tick.date())
        if x < plot_l - Pt(5) or x > plot_l + plot_w + Pt(5):
            continue
        etiqueta = f"{MESES_ES_ABBR[tick.month - 1]}-{str(tick.year)[2:]}"
        forma_texto(slide, x - Pt(16), baseline + Pt(4), Pt(32), Pt(22), etiqueta,
                   size=7.5, color_hex=COL_MUTED, bold=False, align=PP_ALIGN.CENTER,
                   anchor=MSO_ANCHOR.TOP, rotation=45, wrap=False)


# ---------------------------------------------------------------------------
# Placeholders (equivalente a officer::ph_location_label)
# ---------------------------------------------------------------------------
def buscar_layout(prs, nombre: str):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == nombre:
                return layout
    raise ValueError(f"No se encontró el layout: {nombre!r}")


def ph_por_nombre(contenedor, nombre):
    for ph in contenedor.placeholders:
        if ph.name == nombre:
            return ph
    return None


def ph_en_slide_por_nombre(slide, layout, nombre):
    """Busca en `slide` el placeholder correspondiente al `nombre` definido
    en `layout`. python-pptx NO conserva el nombre personalizado del
    placeholder al clonarlo del layout hacia la diapositiva — solo conserva
    el `idx`. Por eso se resuelve el idx en el layout y se busca ese idx en
    la diapositiva, en vez de comparar nombres directamente en la diapositiva."""
    lp = ph_por_nombre(layout, nombre)
    if lp is None:
        return ph_por_nombre(slide, nombre)
    idx = lp.placeholder_format.idx
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def geom_placeholder(slide, layout, nombre):
    """Devuelve (placeholder_en_slide, (L, T, W, H)) con herencia del layout."""
    ph = ph_en_slide_por_nombre(slide, layout, nombre)
    lp = ph_por_nombre(layout, nombre)
    ref = ph if ph is not None else lp
    if ref is None:
        return None, (None, None, None, None)

    def _pick(attr):
        v = getattr(ph, attr) if ph is not None else None
        if v is None and lp is not None:
            v = getattr(lp, attr)
        return v

    return ph, (_pick("left"), _pick("top"), _pick("width"), _pick("height"))


def quitar_placeholder(ph):
    if ph is not None:
        ph._element.getparent().remove(ph._element)


def set_texto_placeholder(slide, layout, nombre, texto):
    ph = ph_en_slide_por_nombre(slide, layout, nombre)
    if ph is not None:
        ph.text = texto


def dibujar_en_placeholder(slide, layout, nombre, funcion_dibujo, *args, **kwargs):
    """Resuelve el placeholder `nombre`, lo quita, y llama a
    `funcion_dibujo(slide, box, *args, **kwargs)` para dibujar formas
    nativas dentro de su geometría (L, T, W, H)."""
    ph, box = geom_placeholder(slide, layout, nombre)
    if box[0] is None:
        return
    quitar_placeholder(ph)
    funcion_dibujo(slide, box, *args, **kwargs)
