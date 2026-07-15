# -*- coding: utf-8 -*-
"""Generador del reporte de productividad en PowerPoint — modo "Plan de Justicia".

Las formas nativas (rectángulos, línea freeform, gráficas de barras y de
serie temporal) y las funciones de manejo de placeholders viven en
`py.formas_nativas`, compartidas con el modo "CLUES". Este archivo solo
tiene lo específico de Plan de Justicia: tarjetas de valor, tablas, y el
armado de datos de `crear_reporte_productividad`.

La función principal `crear_reporte_productividad` devuelve un objeto
`Presentation` de python-pptx. Para guardar:

    prs = crear_reporte_productividad(...)
    prs.save("reporte.pptx")
"""

from __future__ import annotations

from datetime import date, timedelta
from pathlib import Path

import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE

from py.formas_nativas import (
    MESES_ES,
    rgb as _rgb,
    fmt_num,
    buscar_layout,
    geom_placeholder,
    quitar_placeholder,
    set_texto_placeholder,
    dibujar_en_placeholder,
    floor_month,
    ceiling_month,
    mes_titulo,
    mes_abbr_titulo,
    dibujar_grafica_planeacion_historica,
    dibujar_grafica_planeacion_2024_2026,
    dibujar_grafica_consultas_periodos,
)

# Colores ------------------------------------------------------------------
col_rojo_chillon = "#FF0000"
col_amarillo_chillon = "#FFC107"
col_verde_chillon = "#00B050"
# Color de letras y bordes
col_muted = "#6B7280"
col_borde = "#D1D5DB"
col_texto = "#111827"
# Color de gráficas
col_verde = "#1E5B4F"   # IMSS verde
col_guinda = "#611232"  # guinda
col_dorado = "#A57F2C"
col_verde_pastel = "#FFFFFF"


def _es_vacio(x) -> bool:
    if x is None:
        return True
    if isinstance(x, (list, tuple, np.ndarray, pd.Series)):
        return len(x) == 0
    return False


def _primer_valor(x):
    """Equivalente a x[1] de R (primer elemento) para escalares o vectores."""
    if isinstance(x, (list, tuple, np.ndarray, pd.Series)):
        return x[0] if len(x) else np.nan
    return x


# ---------------------------------------------------------------------------
# Funciones de valuebox
# ---------------------------------------------------------------------------
def fmt_delta(x) -> dict:
    """Formatea la variación porcentual: etiqueta, color e ícono."""
    if x is None:
        return {"label": "", "col": col_muted, "icon": ""}

    if isinstance(x, (list, tuple, np.ndarray, pd.Series)):
        if len(x) == 0:
            return {"label": "", "col": col_muted, "icon": ""}
        if pd.isna(pd.Series(list(x))).all():
            return {"label": "s/d", "col": col_muted, "icon": ""}
        x = x[0]

    if pd.isna(x):
        return {"label": "s/d", "col": col_muted, "icon": ""}

    x = _primer_valor(x)

    if x > 0:
        return {"label": f"+{fmt_pct_int(x)}%", "col": col_verde, "icon": "▲ "}
    if x < 0:
        return {"label": f"{fmt_pct_int(x)}%", "col": col_guinda, "icon": "▼ "}
    return {"label": "0%", "col": col_muted, "icon": "• "}


def fmt_pct_int(x) -> str:
    """Representa el porcentaje sin decimales (como R con round(...,0))."""
    return str(int(round(float(x))))


def elige_acento(var_2025, var_2024,
                 verde=col_verde_chillon,
                 amarillo=col_amarillo_chillon,
                 rojo=col_rojo_chillon) -> str:
    """Elige el color de acento según cuántas variaciones son negativas."""
    v25_neg = (var_2025 is not None) and (not pd.isna(var_2025)) and var_2025 < 0
    v24_neg = (var_2024 is not None) and (not pd.isna(var_2024)) and var_2024 < 0

    if v25_neg and v24_neg:
        return rojo
    if v25_neg != v24_neg:  # xor
        return amarillo
    return verde


def crear_card_institucional(numero, titulo, var_vs_2025, var_vs_2024,
                             acento=col_verde,
                             size_num=18, size_titulo=11, size_delta=9) -> dict:
    """Devuelve una especificación de tarjeta (se dibuja al colocarla).

    En R esto devolvía un objeto `rvg::dml`; aquí devolvemos un dict con la
    información necesaria para dibujar la tarjeta como formas nativas dentro
    del placeholder de destino (ver `dibujar_card`).
    """
    d25 = fmt_delta(var_vs_2025)
    d24 = fmt_delta(var_vs_2024)

    mostrar_comparativos = (var_vs_2025 is not None) and (var_vs_2024 is not None)

    return {
        "numero": fmt_num(numero),
        "titulo": titulo,
        "d25": d25,
        "d24": d24,
        "acento": acento,
        "mostrar_comparativos": mostrar_comparativos,
        "size_num": size_num,
        "size_titulo": size_titulo,
        "size_delta": size_delta,
    }


def dibujar_card(slide, box, spec):
    """Dibuja la tarjeta como formas nativas dentro del rectángulo `box`.

    box = (left, top, width, height) en EMU.
    """
    L, T, W, H = box

    # --- Card base (rectángulo redondeado)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, L, T, W, H)
    card.fill.solid()
    card.fill.fore_color.rgb = _rgb("#FFFFFF")
    card.line.color.rgb = _rgb(col_borde)
    card.line.width = Pt(1)
    card.shadow.inherit = False

    # --- Línea lateral (acento)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        L + int(0.025 * W), T + int(0.05 * H),
        int(0.02 * W), int(0.90 * H),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(spec["acento"])
    bar.line.fill.background()
    bar.shadow.inherit = False

    def _txt(top_frac, height_frac):
        tb = slide.shapes.add_textbox(
            L + int(0.07 * W), T + int(top_frac * H),
            int(0.90 * W), int(height_frac * H),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        return tf

    # --- Número
    tf = _txt(0.12, 0.26)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = spec["numero"]
    r.font.bold = True
    r.font.size = Pt(spec["size_num"])
    r.font.color.rgb = _rgb(col_dorado)
    r.font.name = "Calibri"

    # --- Título
    tf = _txt(0.42, 0.18)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = spec["titulo"]
    r.font.bold = True
    r.font.size = Pt(spec["size_titulo"])
    r.font.color.rgb = _rgb(col_texto)
    r.font.name = "Calibri"

    # --- Comparativos
    if spec["mostrar_comparativos"]:
        _linea_comparativo(_txt(0.64, 0.15), spec["d25"], "vs 2025 ", spec["size_delta"])
        _linea_comparativo(_txt(0.79, 0.15), spec["d24"], "vs 2024 ", spec["size_delta"])


def _linea_comparativo(tf, delta, etiqueta, size_delta):
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    r1 = p.add_run()
    r1.text = f"{delta['icon']}{etiqueta}"
    r1.font.size = Pt(size_delta)
    r1.font.color.rgb = _rgb(col_muted)
    r1.font.name = "Calibri"

    r2 = p.add_run()
    r2.text = delta["label"]
    r2.font.size = Pt(size_delta)
    r2.font.bold = True
    r2.font.color.rgb = _rgb(delta["col"])
    r2.font.name = "Calibri"


# ---------------------------------------------------------------------------
# Helpers de datos (dplyr/tidyr -> pandas)
# ---------------------------------------------------------------------------
def rellenar_anios(df: pd.DataFrame, anios) -> pd.DataFrame:
    """Asegura una fila por año del rango, rellenando numéricos con 0."""
    base = pd.DataFrame({"anio": list(anios)})
    out = base.merge(df, on="anio", how="left")
    num_cols = out.select_dtypes(include=[np.number]).columns
    out[num_cols] = out[num_cols].fillna(0)
    return out


def asegurar_columnas(df: pd.DataFrame, cols) -> pd.DataFrame:
    """Agrega columnas faltantes con valor 0."""
    df = df.copy()
    for c in cols:
        if c not in df.columns:
            df[c] = 0
    return df


def obtener_col(df: pd.DataFrame, nombre, default=0):
    """Devuelve la columna si existe; si no, un vector de `default`."""
    if nombre in df.columns:
        return df[nombre]
    return pd.Series([default] * len(df), index=df.index)


def valor_anio_col(df: pd.DataFrame, columna, anio_objetivo):
    """Primer valor de `columna` para el año dado (NaN si no existe)."""
    if columna not in df.columns:
        return np.nan
    val = df.loc[df["anio"] == anio_objetivo, columna]
    if len(val) == 0:
        return np.nan
    return val.iloc[0]


def tiene_dato_2026(df: pd.DataFrame, columna) -> bool:
    if columna not in df.columns:
        return False
    valor = df.loc[df["anio"] == 2026, columna]
    if len(valor) == 0:
        return False
    return bool(((~valor.isna()) & (valor != 0)).any())


def hay_indicador_2026(df: pd.DataFrame, columna) -> bool:
    val = valor_anio_col(df, columna, 2026)
    return (not pd.isna(val)) and val != 0


# ---------------------------------------------------------------------------
# Value boxes
# ---------------------------------------------------------------------------
def crear_valueboxes_2026(df_3anios, mapa_titulos, sufijo="",
                          incluir_comparativos=True,
                          acento_sin_comparativo="#B0B0B0",
                          size_num=18, size_titulo=11, size_delta=9) -> dict:
    """Crea un dict {clave+sufijo: spec_de_tarjeta} para el año 2026."""
    d26 = df_3anios[df_3anios["anio"] == 2026]

    cards = {}
    for var, titulo in mapa_titulos.items():
        numero = d26[var] if var in d26.columns else pd.Series(dtype=float)

        if len(numero) == 0 or numero.isna().all():
            numero = 0
        else:
            numero = numero.iloc[0]

        if incluir_comparativos:
            col25 = f"var_2026_vs_2025_{var}"
            col24 = f"var_2026_vs_2024_{var}"

            var_vs_2025 = d26[col25].iloc[0] if col25 in d26.columns and len(d26) else np.nan
            var_vs_2024 = d26[col24].iloc[0] if col24 in d26.columns and len(d26) else np.nan

            acento = elige_acento(var_vs_2025, var_vs_2024)
        else:
            var_vs_2025 = None
            var_vs_2024 = None
            acento = acento_sin_comparativo

        cards[f"{var}{sufijo}"] = crear_card_institucional(
            numero=numero, titulo=titulo,
            var_vs_2025=var_vs_2025, var_vs_2024=var_vs_2024,
            acento=acento,
            size_num=size_num, size_titulo=size_titulo, size_delta=size_delta,
        )
    return cards


def definir_layout_valueboxes(datos_consulta_funcion) -> dict:
    """Elige el layout y las métricas según qué indicadores tienen dato 2026."""
    hay_general = tiene_dato_2026(datos_consulta_funcion, "consulta_gral")
    hay_esp = tiene_dato_2026(datos_consulta_funcion, "consulta_esp")
    hay_qx = tiene_dato_2026(datos_consulta_funcion, "qx")
    hay_egresos = tiene_dato_2026(datos_consulta_funcion, "egresos")

    if hay_general and not hay_esp and not hay_qx and not hay_egresos:
        return {"layout": "2_valueboxes", "metricas": ["consulta_gral"]}

    if hay_general and hay_esp and not hay_qx and not hay_egresos:
        return {"layout": "6_valueboxes",
                "metricas": ["total_consultas", "consulta_gral", "consulta_esp"]}

    if hay_general and hay_esp and hay_qx and not hay_egresos:
        return {"layout": "8_valueboxes",
                "metricas": ["total_consultas", "consulta_gral", "consulta_esp", "qx"]}

    if hay_general and hay_esp and hay_qx and hay_egresos:
        return {"layout": "10_valueboxes",
                "metricas": ["total_consultas", "consulta_gral", "consulta_esp", "qx", "egresos"]}

    # Fallback por combinaciones inusuales
    metricas_presentes = []
    if hay_general:
        metricas_presentes.append("consulta_gral")
    if hay_esp:
        metricas_presentes.append("consulta_esp")
    if hay_qx:
        metricas_presentes.append("qx")
    if hay_egresos:
        metricas_presentes.append("egresos")

    if len(metricas_presentes) >= 2:
        metricas_presentes = ["total_consultas"] + metricas_presentes

    n = len(metricas_presentes)
    if n <= 1:
        layout_fallback = "2_valueboxes"
    elif n == 3:
        layout_fallback = "6_valueboxes"
    elif n == 4:
        layout_fallback = "8_valueboxes"
    else:  # n >= 5
        layout_fallback = "10_valueboxes"

    return {"layout": layout_fallback, "metricas": metricas_presentes}


def definir_layout_historico(datos_consulta_funcion):
    hay_consultas = any([
        hay_indicador_2026(datos_consulta_funcion, "total_consultas"),
        hay_indicador_2026(datos_consulta_funcion, "consulta_gral"),
        hay_indicador_2026(datos_consulta_funcion, "consulta_esp"),
    ])
    hay_proc = any([
        hay_indicador_2026(datos_consulta_funcion, "qx"),
        hay_indicador_2026(datos_consulta_funcion, "egresos"),
    ])

    if hay_consultas and hay_proc:
        return "Historico consultas y procedimientos"
    if hay_consultas:
        return "Historico consultas"
    return None


# ---------------------------------------------------------------------------
# Colocación de tarjetas en el .pptx (los helpers de placeholder genéricos
# — buscar_layout, geom_placeholder, etc. — están en py.formas_nativas)
# ---------------------------------------------------------------------------
def _colocar_card(slide, layout, nombre, spec):
    ph, box = geom_placeholder(slide, layout, nombre)
    if box[0] is None:
        return
    quitar_placeholder(ph)
    dibujar_card(slide, box, spec)


def imprimir_valueboxes_dinamicos(prs, layout_name, boxes_superior, boxes_inferior,
                                  titulo="Productividad IMSS Bienestar",
                                  fecha=None):
    """Agrega una diapositiva de value boxes y coloca las tarjetas."""
    layout = buscar_layout(prs, layout_name)
    slide = prs.slides.add_slide(layout)

    set_texto_placeholder(slide, layout, "Título 1", titulo)
    if fecha is not None:
        set_texto_placeholder(slide, layout, "fecha", fecha)

    for i, spec in enumerate(boxes_superior, start=1):
        if spec is not None:
            _colocar_card(slide, layout, f"arriba {i}", spec)

    for i, spec in enumerate(boxes_inferior, start=1):
        if spec is not None:
            _colocar_card(slide, layout, f"abajo {i}", spec)

    return slide



# ---------------------------------------------------------------------------
# Tablas (flextable -> tabla nativa de python-pptx)
# ---------------------------------------------------------------------------
def armar_tabla_dinamica(df, indicadores, etiquetas, mes_nombre) -> pd.DataFrame:
    """Construye la tabla comparativa 2025 vs 2026 con crecimiento anual."""
    valor_2025 = [valor_anio_col(df, ind, 2025) for ind in indicadores]
    valor_2026 = [valor_anio_col(df, ind, 2026) for ind in indicadores]

    filas = []
    for etq, v25, v26 in zip(etiquetas, valor_2025, valor_2026):
        if pd.isna(v25) or v25 == 0:
            crecimiento = np.nan
        else:
            crecimiento = round((v26 - v25) / v25 * 100, 0)

        if pd.isna(crecimiento):
            crec_txt = "s/d"
        elif crecimiento > 0:
            crec_txt = f"+{int(crecimiento)} %"
        else:
            crec_txt = f"{int(crecimiento)} %"

        v25_txt = "s/d" if pd.isna(v25) else fmt_num(v25)
        v26_txt = "s/d" if pd.isna(v26) else fmt_num(v26)

        filas.append({
            "indicador": etq,
            f"{mes_nombre} 2025": v25_txt,
            f"{mes_nombre} 2026": v26_txt,
            "Crecimiento anual": crec_txt,
        })

    return pd.DataFrame(filas, columns=[
        "indicador", f"{mes_nombre} 2025", f"{mes_nombre} 2026", "Crecimiento anual"])


def _set_cell_border(cell, color="6B7280", width_pt=1):
    """Agrega borde a las 4 aristas de una celda (python-pptx no lo expone)."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for edge in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        tag = qn(edge)
        for existing in tcPr.findall(tag):
            tcPr.remove(existing)
        ln = tcPr.makeelement(tag, {
            "w": str(int(Pt(width_pt))), "cap": "flat",
            "cmpd": "sng", "algn": "ctr",
        })
        fill = ln.makeelement(qn("a:solidFill"), {})
        clr = fill.makeelement(qn("a:srgbClr"), {"val": color})
        fill.append(clr)
        ln.append(fill)
        tcPr.append(ln)


def ft_planeacion(prs_slide, layout, nombre_ph, df,
                  w1=4.05, w2=1.20, w3=1.20, w4=1.35,
                  header_negro="#3B3B3B", verde="#2F6F63", menta="#D9F2EE",
                  size_header=None, size_body=None, h_fila=None):
    """Dibuja una tabla estilizada en el placeholder indicado."""
    n_filas = len(df)

    if size_header is None or size_body is None or h_fila is None:
        if n_filas == 1:
            size_header, size_body, h_fila = 13, 12, 0.45
        elif n_filas == 2:
            size_header, size_body, h_fila = 11, 10, 0.35
        else:
            size_header, size_body, h_fila = 10, 9, 0.28

    ph, (L, T, W, H) = geom_placeholder(prs_slide, layout, nombre_ph)
    if L is None:
        return
    quitar_placeholder(ph)

    encabezados = ["Indicador"] + list(df.columns[1:])
    n_cols = len(encabezados)
    anchos = [w1, w2, w3, w4][:n_cols]
    ancho_total = Inches(sum(anchos))
    alto_total = Inches(h_fila * (n_filas + 1))

    gf = prs_slide.shapes.add_table(n_filas + 1, n_cols, L, T, ancho_total, alto_total)
    tabla = gf.table
    tabla.first_row = False
    tabla.horz_banding = False

    for j, wi in enumerate(anchos):
        tabla.columns[j].width = Inches(wi)
    for i in range(n_filas + 1):
        tabla.rows[i].height = Inches(h_fila)

    # Encabezado
    for j, texto in enumerate(encabezados):
        cell = tabla.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(header_negro)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(texto)
        r.font.bold = True
        r.font.size = Pt(size_header)
        r.font.color.rgb = _rgb("#FFFFFF")
        _set_cell_border(cell)

    # Cuerpo
    for i in range(n_filas):
        for j, colname in enumerate(df.columns):
            cell = tabla.cell(i + 1, j)
            if j == 2:  # tercera columna (año 2026) -> fondo menta
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(menta)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb("#FFFFFF")
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(df.iloc[i, j])
            r.font.size = Pt(size_body)
            r.font.color.rgb = _rgb(col_texto)
            _set_cell_border(cell)

    return tabla


def _calcular_fecha_corte(hoy: date) -> date:
    """Reproduce la lógica de corte semanal (miércoles) del código R."""
    # POSIXlt$wday: domingo=0 ... sábado=6
    wday0 = (hoy.weekday() + 1) % 7
    if wday0 == 3:  # miércoles (lubridate::wday == 4)
        return hoy - timedelta(days=7)
    return hoy - timedelta(days=((wday0 + 4) % 7))


# ---------------------------------------------------------------------------
# FUNCIÓN PRINCIPAL
# ---------------------------------------------------------------------------
_RUTA_MASTER_DEFAULT = (
    Path(__file__).resolve().parent.parent.parent / "data" / "master_presentacion.pptx"
)


def crear_reporte_productividad(codigo_clues, clues_info, metas, historicos,
                                procedimientos_personas=None,
                                ruta_master=None,
                                hoy=None):
    """Genera la presentación de productividad y devuelve el objeto Presentation."""
    if hoy is None:
        hoy = date.today()
    if ruta_master is None:
        ruta_master = str(_RUTA_MASTER_DEFAULT)

    # Variables de corte ----------------------------------------------------
    fecha_corte = _calcular_fecha_corte(hoy)
    fecha_fin_graf = floor_month(fecha_corte)

    fecha_portada = f"{fecha_corte.day:02d} de {MESES_ES[fecha_corte.month - 1]} de {fecha_corte.year}"
    mes_nombre = MESES_ES[fecha_corte.month - 1].capitalize()

    # Filtrado por plan / nacional -----------------------------------------
    if codigo_clues == "NACIONAL":
        clues_info_filtrado = pd.DataFrame(
            [{"clues": "NACIONAL", "nombre": "NACIONAL", "entidad": "NACIONAL"}])
        metas_filtrado = metas.copy()
    else:
        mask = (clues_info["clues"] == codigo_clues)
        if "id" in clues_info.columns:
            mask = mask | (clues_info["id"] == codigo_clues)
        clues_info_filtrado = clues_info[mask].copy()
        metas_filtrado = metas[metas["clues"] == codigo_clues].copy()

    if len(metas_filtrado) == 0:
        raise ValueError(f"No hay metas para: {codigo_clues}")

    if len(clues_info_filtrado) == 0:
        clues_info_filtrado = pd.DataFrame(
            [{"clues": codigo_clues, "nombre": codigo_clues, "entidad": np.nan}])

    meta_total_consultas = (
        pd.to_numeric(metas_filtrado.get("meta_general_anual"), errors="coerce").fillna(0).sum()
        + pd.to_numeric(metas_filtrado.get("meta_especialidad_anual"), errors="coerce").fillna(0).sum()
    )
    meta_qx = pd.to_numeric(metas_filtrado.get("meta_cirugia_anual"), errors="coerce").fillna(0).sum()

    # Derivar procedimientos_personas si no viene ---------------------------
    if procedimientos_personas is None:
        h = historicos.copy()
        h["anio"] = pd.to_datetime(h["fecha"]).dt.year
        agg = h.groupby("anio", as_index=False).agg(
            total_consultas=("consulta_total", "sum"),
            consulta_gral=("consulta_general", "sum"),
            consulta_esp=("consulta_especialidad", "sum"),
            qx=("procedimientos_qx", "sum"),
            egresos=("egresos", "sum"),
        )
        largo = agg.melt(id_vars="anio",
                         value_vars=["total_consultas", "consulta_gral",
                                     "consulta_esp", "qx", "egresos"],
                         var_name="tipo_procedimiento", value_name="procedimientos")
        largo["fecha"] = largo["anio"]
        largo["personas"] = largo["procedimientos"]
        procedimientos_personas = largo[["fecha", "tipo_procedimiento",
                                         "procedimientos", "personas"]]

    procedimientos_personas = procedimientos_personas.copy()
    recodificar = {
        "consulta total": "total_consultas",
        "general": "consulta_gral",
        "especialidad": "consulta_esp",
        "qx": "qx",
        "egresos": "egresos",
    }
    procedimientos_personas["tipo_procedimiento"] = (
        procedimientos_personas["tipo_procedimiento"].map(lambda v: recodificar.get(v, v)))

    # Presentación ----------------------------------------------------------
    prs = Presentation(ruta_master)

    # Portada ---------------------------------------------------------------
    layout_portada = buscar_layout(prs, "Portada 3")
    slide = prs.slides.add_slide(layout_portada)
    titulo_portada = f"Reporte de {str(clues_info_filtrado['clues'].iloc[0]).title()}"
    set_texto_placeholder(slide, layout_portada, "Título 1", titulo_portada)
    set_texto_placeholder(slide, layout_portada, "Marcador de contenido 2", fecha_portada)

    # Datos base ------------------------------------------------------------
    cols_metricas = ["consulta_gral", "consulta_esp", "qx", "total_consultas", "egresos"]

    h = historicos.copy()
    h["anio"] = pd.to_datetime(h["fecha"]).dt.year
    # obtener_col: columnas faltantes se tratan como 0
    h["_cgral"] = obtener_col(h, "consulta_general")
    h["_cesp"] = obtener_col(h, "consulta_especialidad")
    h["_qx"] = obtener_col(h, "procedimientos_qx")
    h["_tot"] = obtener_col(h, "consulta_total")
    h["_egr"] = obtener_col(h, "egresos")
    datos_anual = h.groupby("anio", as_index=False).agg(
        consulta_gral_anual=("_cgral", "sum"),
        consulta_esp_anual=("_cesp", "sum"),
        qx_anual=("_qx", "sum"),
        total_consultas_anual=("_tot", "sum"),
        egresos_anual=("_egr", "sum"),
    )
    datos_anual = rellenar_anios(datos_anual, range(2020, 2027))

    # datos_consulta_funcion (procedimientos)
    dcf = procedimientos_personas.copy()
    dcf["anio"] = pd.to_numeric(dcf["fecha"], errors="coerce")
    dcf = (dcf[["anio", "tipo_procedimiento", "procedimientos"]]
           .pivot_table(index="anio", columns="tipo_procedimiento",
                        values="procedimientos", aggfunc="sum")
           .reset_index())
    dcf.columns.name = None
    dcf = rellenar_anios(dcf, range(2024, 2027))
    dcf = asegurar_columnas(dcf, cols_metricas)
    datos_consulta_funcion = dcf.merge(datos_anual, on="anio", how="left").sort_values("anio")
    num_cols = datos_consulta_funcion.select_dtypes(include=[np.number]).columns
    datos_consulta_funcion[num_cols] = datos_consulta_funcion[num_cols].fillna(0)
    datos_consulta_funcion["total_consultas_meta"] = np.where(
        datos_consulta_funcion["anio"] == 2026,
        meta_total_consultas, datos_consulta_funcion["total_consultas_anual"])
    datos_consulta_funcion["qx_meta"] = np.where(
        datos_consulta_funcion["anio"] == 2026,
        meta_qx, datos_consulta_funcion["qx_anual"])

    # datos_curps (personas)
    dcu = procedimientos_personas.copy()
    dcu["anio"] = pd.to_numeric(dcu["fecha"], errors="coerce")
    dcu = (dcu[["anio", "tipo_procedimiento", "personas"]]
           .pivot_table(index="anio", columns="tipo_procedimiento",
                        values="personas", aggfunc="sum")
           .reset_index())
    dcu.columns.name = None
    dcu = rellenar_anios(dcu, range(2024, 2027))
    dcu = asegurar_columnas(dcu, cols_metricas)
    datos_curps = dcu
    num_cols = datos_curps.select_dtypes(include=[np.number]).columns
    datos_curps[num_cols] = datos_curps[num_cols].fillna(0)

    # Comparativos ----------------------------------------------------------
    for col in cols_metricas:
        for ref in (2024, 2025):
            nombre_var = f"var_2026_vs_{ref}_{col}"

            v26 = valor_anio_col(datos_consulta_funcion, col, 2026)
            vref = valor_anio_col(datos_consulta_funcion, col, ref)
            cond = (not pd.isna(v26)) and (not pd.isna(vref)) and vref != 0
            valor = round(100 * (1 - (vref / v26)), 0) if cond and v26 not in (0, np.nan) else 0
            datos_consulta_funcion[nombre_var] = np.where(
                datos_consulta_funcion["anio"] == 2026, valor, 0)

            v26c = valor_anio_col(datos_curps, col, 2026)
            vrefc = valor_anio_col(datos_curps, col, ref)
            condc = (not pd.isna(v26c)) and (not pd.isna(vrefc)) and vrefc != 0
            valorc = round(100 * (1 - (vrefc / v26c)), 0) if condc and v26c not in (0, np.nan) else 0
            datos_curps[nombre_var] = np.where(datos_curps["anio"] == 2026, valorc, 0)

    # Value boxes -----------------------------------------------------------
    mapa_titulos_consultas = {
        "total_consultas": "Consultas totales",
        "consulta_gral": "Consulta general",
        "consulta_esp": "Especialidad",
        "qx": "Procedimientos quirúrgicos",
        "egresos": "Egresos",
    }
    mapa_titulos_curp = {
        "total_consultas": "Consultas totales",
        "consulta_gral": "Consulta general",
        "consulta_esp": "Especialidad",
        "qx": "Intervenidas",
        "egresos": "Egresadas",
    }

    vbox_consultas = crear_valueboxes_2026(
        datos_consulta_funcion, mapa_titulos_consultas, incluir_comparativos=True)
    vbox_curps = crear_valueboxes_2026(
        datos_curps, mapa_titulos_curp, sufijo="_p", incluir_comparativos=True)

    config_vb = definir_layout_valueboxes(datos_consulta_funcion)
    layout_vb = config_vb["layout"]
    metricas_vb = config_vb["metricas"]

    lista_boxes_superior = {
        "total_consultas": vbox_consultas.get("total_consultas"),
        "consulta_gral": vbox_consultas.get("consulta_gral"),
        "consulta_esp": vbox_consultas.get("consulta_esp"),
        "qx": vbox_consultas.get("qx"),
        "egresos": vbox_consultas.get("egresos"),
    }
    lista_boxes_inferior = {
        "total_consultas": vbox_curps.get("total_consultas_p"),
        "consulta_gral": vbox_curps.get("consulta_gral_p"),
        "consulta_esp": vbox_curps.get("consulta_esp_p"),
        "qx": vbox_curps.get("qx_p"),
        "egresos": vbox_curps.get("egresos_p"),
    }

    imprimir_valueboxes_dinamicos(
        prs=prs,
        layout_name=layout_vb,
        boxes_superior=[lista_boxes_superior[m] for m in metricas_vb],
        boxes_inferior=[lista_boxes_inferior[m] for m in metricas_vb],
        titulo="Productividad IMSS Bienestar",
        fecha=f"Del 01 de enero al {fecha_portada}",
    )

    # Diapo 3: histórico 2020-2025 -----------------------------------------
    fecha_corte_15 = date(fecha_corte.year, fecha_corte.month, 15)
    corte_md = f"{fecha_corte_15.month:02d}-{fecha_corte_15.day:02d}"

    hh = historicos.copy()
    hh["fecha_dt"] = pd.to_datetime(hh["fecha"]).dt.tz_localize(None)
    hh["consulta_general_tmp"] = obtener_col(hh, "consulta_general") + obtener_col(hh, "consulta_gral")
    hh["consulta_esp_tmp"] = obtener_col(hh, "consulta_especialidad") + obtener_col(hh, "consulta_esp")
    hh["procedimientos_qx_tmp"] = obtener_col(hh, "procedimientos_qx")
    hh["egresos_tmp"] = obtener_col(hh, "egresos")
    hh["md"] = hh["fecha_dt"].dt.strftime("%m-%d")
    hh = hh[(hh["fecha_dt"].dt.year >= 2020) &
            (hh["fecha_dt"].dt.year <= 2025) &
            (hh["md"] <= corte_md)].copy()
    hh["anio"] = hh["fecha_dt"].dt.year.astype(str)

    datos_historicos_2020_2025 = hh.groupby("anio", as_index=False).agg(
        consulta_gral=("consulta_general_tmp", "sum"),
        consulta_esp=("consulta_esp_tmp", "sum"),
        qx=("procedimientos_qx_tmp", "sum"),
        egresos=("egresos_tmp", "sum"),
    )

    def _total_consultas(g, e):
        if g > 0 and e > 0:
            return g + e
        if g > 0:
            return g
        if e > 0:
            return e
        return 0

    datos_historicos_2020_2025["total_consultas"] = [
        _total_consultas(g, e) for g, e in
        zip(datos_historicos_2020_2025["consulta_gral"],
            datos_historicos_2020_2025["consulta_esp"])]

    da_str = datos_anual.copy()
    da_str["anio"] = da_str["anio"].astype(str)
    datos_historicos_2020_2025 = datos_historicos_2020_2025.merge(
        da_str, on="anio", how="left")
    for c in ["consulta_gral_anual", "consulta_esp_anual", "qx_anual"]:
        datos_historicos_2020_2025[c] = datos_historicos_2020_2025.get(c, 0)
        datos_historicos_2020_2025[c] = datos_historicos_2020_2025[c].fillna(0)
    datos_historicos_2020_2025["total_consultas_anual"] = [
        _total_consultas(g, e) for g, e in
        zip(datos_historicos_2020_2025["consulta_gral_anual"],
            datos_historicos_2020_2025["consulta_esp_anual"])]
    datos_historicos_2020_2025 = datos_historicos_2020_2025.sort_values("anio")

    print("REVISION DATOS HISTORICOS")
    print(datos_historicos_2020_2025)

    hay_consultas_2020 = bool((datos_historicos_2020_2025["total_consultas"] > 0).any())
    hay_qx_2020 = bool(((datos_historicos_2020_2025["qx"] > 0) |
                        (datos_historicos_2020_2025["egresos"] > 0)).any())

    if hay_consultas_2020:
        if hay_qx_2020:
            lay = buscar_layout(prs, "1_Historico consultas y procedimientos")
            s = prs.slides.add_slide(lay)
            set_texto_placeholder(s, lay, "Título 1", "Productividad IMSS Bienestar")
            dibujar_en_placeholder(s, lay, "Grafica 1", dibujar_grafica_planeacion_historica,
                                    datos_historicos_2020_2025, "total_consultas_anual",
                                    "total_consultas", "Consultas totales")
            dibujar_en_placeholder(s, lay, "Grafica 2", dibujar_grafica_planeacion_historica,
                                    datos_historicos_2020_2025, "qx_anual", "qx",
                                    "Procedimientos quirúrgicos")
            set_texto_placeholder(s, lay, "fecha", f"Del 01 de enero al {fecha_portada}")
        else:
            lay = buscar_layout(prs, "1_Historico consultas")
            s = prs.slides.add_slide(lay)
            set_texto_placeholder(s, lay, "Título 1", "Productividad IMSS Bienestar")
            dibujar_en_placeholder(s, lay, "Grafica 1", dibujar_grafica_planeacion_historica,
                                    datos_historicos_2020_2025, "total_consultas_anual",
                                    "total_consultas", "Consultas totales")
            set_texto_placeholder(s, lay, "fecha", f"Del 01 de enero al {fecha_portada}")

    # Diapo 4: 2024-2026 con tablas ----------------------------------------
    for col in ("consulta_gral", "consulta_esp", "qx", "egresos"):
        if col not in datos_consulta_funcion.columns:
            datos_consulta_funcion[col] = 0

    d2426 = datos_consulta_funcion.copy()
    d2426["anio_num"] = d2426["anio"].astype(int)
    for c in ("consulta_gral", "consulta_esp", "qx", "egresos",
              "consulta_gral_anual", "consulta_esp_anual", "qx_anual"):
        if c in d2426.columns:
            d2426[c] = d2426[c].fillna(0)

    d2426["total_consultas"] = [
        _total_consultas(g, e) for g, e in
        zip(d2426["consulta_gral"], d2426["consulta_esp"])]

    def _total_meta(row):
        if row["anio_num"] == 2026 and not pd.isna(row.get("total_consultas_meta")):
            return row["total_consultas_meta"]
        if row["anio_num"] == 2026:
            return row["total_consultas"]
        g, e = row.get("consulta_gral_anual", 0), row.get("consulta_esp_anual", 0)
        if g > 0 and e > 0:
            return g + e
        if g > 0:
            return g
        if e > 0:
            return e
        if row.get("total_consultas_anual", 0) > 0:
            return row["total_consultas_anual"]
        return row["total_consultas"]

    def _qx_meta(row):
        if row["anio_num"] == 2026 and not pd.isna(row.get("qx_meta")):
            return row["qx_meta"]
        if row["anio_num"] == 2026:
            return row["qx"]
        if row.get("qx_anual", 0) > 0:
            return row["qx_anual"]
        return row["qx"]

    d2426["total_consultas_meta"] = d2426.apply(_total_meta, axis=1)
    d2426["qx_meta"] = d2426.apply(_qx_meta, axis=1)
    d2426 = d2426[d2426["anio_num"].isin([2024, 2025, 2026])]

    print("REVISION DATOS 2024_2026")
    print(d2426)

    hay_consultas_2426 = bool((d2426["total_consultas"] > 0).any())
    hay_qx_2426 = bool(((d2426["qx"] > 0) | (d2426["egresos"] > 0)).any())

    if hay_consultas_2426:
        indicadores_consulta, etiquetas_consulta = [], []
        if hay_indicador_2026(d2426, "consulta_gral"):
            indicadores_consulta.append("consulta_gral")
            etiquetas_consulta.append("Consultas generales")
        if hay_indicador_2026(d2426, "consulta_esp"):
            indicadores_consulta.append("consulta_esp")
            etiquetas_consulta.append("Consultas de especialidad*")

        tabla_consultas = armar_tabla_dinamica(
            d2426, indicadores_consulta, etiquetas_consulta, "Acumulado")

        if hay_qx_2426:
            indicadores_proc, etiquetas_proc = [], []
            if hay_indicador_2026(d2426, "qx"):
                indicadores_proc.append("qx")
                etiquetas_proc.append("Procedimientos quirúrgicos")
            if hay_indicador_2026(d2426, "egresos"):
                indicadores_proc.append("egresos")
                etiquetas_proc.append("Egresos")

            tabla_proc = armar_tabla_dinamica(
                d2426, indicadores_proc, etiquetas_proc, "Acumulado")

            lay = buscar_layout(prs, "Historico consultas y procedimientos")
            s = prs.slides.add_slide(lay)
            set_texto_placeholder(s, lay, "Título 1", "Productividad IMSS Bienestar")
            dibujar_en_placeholder(s, lay, "Grafica 1", dibujar_grafica_planeacion_2024_2026,
                                    d2426, "total_consultas_meta", "total_consultas",
                                    "Consultas totales")
            dibujar_en_placeholder(s, lay, "Grafica 2", dibujar_grafica_planeacion_2024_2026,
                                    d2426, "qx_meta", "qx", "Procedimientos quirúrgicos")
            ft_planeacion(s, lay, "tabla_1", tabla_consultas,
                          w1=2.70, w2=0.90, w3=0.90, w4=0.80,
                          size_header=8, size_body=7.5, h_fila=0.28)
            ft_planeacion(s, lay, "tabla_2", tabla_proc,
                          w1=2.70, w2=0.90, w3=0.90, w4=0.80,
                          size_header=8, size_body=7.5, h_fila=0.28)
            set_texto_placeholder(s, lay, "fecha", f"Del 01 de enero al {fecha_portada}")
        else:
            lay = buscar_layout(prs, "Historico consultas")
            s = prs.slides.add_slide(lay)
            set_texto_placeholder(s, lay, "Título 1", "Productividad IMSS Bienestar")
            dibujar_en_placeholder(s, lay, "Grafica 1", dibujar_grafica_planeacion_2024_2026,
                                    d2426, "total_consultas_meta", "total_consultas",
                                    "Consultas totales")
            ft_planeacion(s, lay, "tabla_1", tabla_consultas,
                          w1=4.60, w2=1.35, w3=1.35, w4=1.40,
                          size_header=11, size_body=10, h_fila=0.38)
            set_texto_placeholder(s, lay, "fecha", f"Del 01 de enero al {fecha_portada}")

    # Diapo 5: serie mensual de consultas ----------------------------------
    sm = historicos.copy()
    sm["fecha"] = pd.to_datetime(sm["fecha"]).dt.tz_localize(None)
    sm["fecha"] = sm["fecha"].apply(lambda d: pd.Timestamp(floor_month(d)))
    sm = sm[sm["fecha"].notna() & sm["consulta_total"].notna()]
    serie_mensual_consultas = (sm.groupby("fecha", as_index=False)
                               .agg(consultas_totales=("consulta_total", "sum"))
                               .sort_values("fecha"))

    subtitulo_5 = f"Agosto 2022 – {MESES_ES[fecha_fin_graf.month - 1].capitalize()} {fecha_fin_graf.year}"

    lay = buscar_layout(prs, "Una grafica")
    s = prs.slides.add_slide(lay)
    set_texto_placeholder(s, lay, "Título 1", "Consultas totales por mes")
    set_texto_placeholder(s, lay, "fecha", subtitulo_5)
    dibujar_en_placeholder(s, lay, "ft", dibujar_grafica_consultas_periodos,
                            serie_mensual_consultas, fecha_inicio="2022-08-01",
                            fecha_fin=str(fecha_fin_graf))

    # Diapo 6: serie mensual de procedimientos quirúrgicos -----------------
    if hay_indicador_2026(datos_consulta_funcion, "qx"):
        smq = historicos.copy()
        smq["fecha"] = pd.to_datetime(smq["fecha"]).dt.tz_localize(None)
        smq["fecha"] = smq["fecha"].apply(lambda d: pd.Timestamp(floor_month(d)))
        smq = smq[smq["fecha"].notna() & smq["consulta_total"].notna()]
        serie_mensual_pq = (smq.groupby("fecha", as_index=False)
                            .agg(consultas_totales=("procedimientos_qx", "sum"))
                            .sort_values("fecha"))

        subtitulo_6 = f"Agosto 2022 – {MESES_ES[fecha_fin_graf.month - 1].capitalize()} {fecha_fin_graf.year}"

        lay = buscar_layout(prs, "Una grafica")
        s = prs.slides.add_slide(lay)
        set_texto_placeholder(s, lay, "Título 1", "Procedimientos quirúrgicos por mes")
        set_texto_placeholder(s, lay, "fecha", subtitulo_6)
        dibujar_en_placeholder(s, lay, "ft", dibujar_grafica_consultas_periodos,
                                serie_mensual_pq, fecha_inicio="2022-08-01",
                                fecha_fin=str(fecha_fin_graf))

    return prs
