"""Generador del reporte PPTX de productividad — modo "CLUES / Unidad médica".

Las formas nativas y las gráficas (barras, serie temporal) viven en
`py.formas_nativas`, compartidas con el modo "Plan de Justicia". Este
archivo tiene lo específico de CLUES: tarjetas de valor, tablas, y el
armado de datos de `crear_reporte_productividad`.
"""

import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from py.clues.charts import fecha_larga_es, mes_anio_es
from py.clues.data_loader import CUBOS_PARQUET, load_categoria_gerencial_nueva
from py.clues.utils_comunes import calcular_fecha_corte
from py.formas_nativas import (
    rgb as _rgb,
    fmt_num,
    dibujar_grafica_planeacion_historica,
    dibujar_grafica_planeacion_2024_2026,
    dibujar_grafica_consultas_periodos,
)

# ---------------------------------------------------------------------------
# Colores
# ---------------------------------------------------------------------------
COL_ROJO_CHILLON = "FF0000"
COL_AMARILLO_CHILLON = "FFC107"
COL_VERDE_CHILLON = "00B050"
COL_MUTED = "6B7280"
COL_BORDE = "D1D5DB"
COL_TEXTO = "111827"
COL_VERDE = "1E5B4F"
COL_GUINDA = "611232"
COL_DORADO = "A57F2C"

_TIPO_PROC_PPTX_MAP = {
    "consulta total": "total_consultas",
    "general": "consulta_gral",
    "especialidad": "consulta_esp",
    "qx": "qx",
    "egresos": "egresos",
}

MAPA_TITULOS_CONSULTAS = {
    "total_consultas": "Consultas totales",
    "consulta_gral": "Consulta general",
    "consulta_esp": "Especialidad",
    "qx": "Procedimientos quirúrgicos",
    "egresos": "Egresos",
}

MAPA_TITULOS_CURP = {
    "total_consultas": "Consultas totales",
    "consulta_gral": "Consulta general",
    "consulta_esp": "Especialidad",
    "qx": "Intervenidas",
    "egresos": "Egresadas",
}


# ---------------------------------------------------------------------------
# Helpers de placeholders (equivalente a officer::ph_location_label)
# ---------------------------------------------------------------------------

def _layout_idx_map(layout) -> dict:
    return {ph.name: ph.placeholder_format.idx for ph in layout.placeholders}


def find_placeholder(layout, slide, name):
    idx_map = _layout_idx_map(layout)
    if name not in idx_map:
        return None
    return slide.placeholders[idx_map[name]]


def _bbox(ph):
    return ph.left, ph.top, ph.width, ph.height


def _remove(shape):
    shape._element.getparent().remove(shape._element)


def add_slide(prs: Presentation, master, layout_name: str):
    layout = next(l for l in master.slide_layouts if l.name == layout_name)
    slide = prs.slides.add_slide(layout)
    return slide, layout


def set_placeholder_text(layout, slide, name: str, text: str):
    ph = find_placeholder(layout, slide, name)
    if ph is None:
        return
    p = ph.text_frame.paragraphs[0]
    if p.runs:
        run = p.runs[0]
        for extra in list(p.runs[1:]):
            extra._r.getparent().remove(extra._r)
    else:
        run = p.add_run()
    run.text = str(text)




def place_native_chart(slide, layout, name: str, funcion_dibujo, *args, **kwargs):
    """Resuelve el placeholder `name` y llama a
    `funcion_dibujo(slide, box, *args, **kwargs)` para dibujar formas nativas
    (editables en PowerPoint) dentro de su geometría, en vez de pegar una
    imagen (como hacía la vieja place_picture)."""
    ph = find_placeholder(layout, slide, name)
    if ph is None:
        return
    box = _bbox(ph)
    _remove(ph)
    funcion_dibujo(slide, box, *args, **kwargs)


def _set_cell_border(cell, color: str = COL_MUTED, width_pt: float = 1.0):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    line_w = str(int(Pt(width_pt)))
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        existing = tcPr.find(qn(tag))
        if existing is not None:
            tcPr.remove(existing)
        ln = tcPr.makeelement(qn(tag), {"w": line_w, "cap": "flat"})
        solid_fill = ln.makeelement(qn("a:solidFill"), {})
        srgb = solid_fill.makeelement(qn("a:srgbClr"), {"val": color})
        solid_fill.append(srgb)
        ln.append(solid_fill)
        tcPr.append(ln)


def place_table(
    slide, layout, name: str, df: pd.DataFrame,
    w: list | None = None,
    header_negro: str = "3B3B3B",
    menta: str = "D9F2EE",
    size_header: int = 10,
    size_body: int = 9,
):
    ph = find_placeholder(layout, slide, name)
    if ph is None:
        return
    left, top, width, height = _bbox(ph)
    _remove(ph)

    n_rows, n_cols = df.shape[0] + 1, df.shape[1]
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    for j, colname in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = "" if colname == "indicador" else str(colname)
        if colname == "indicador":
            cell.text_frame.paragraphs[0].add_run().text = "Indicador"
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(header_negro)
        _set_cell_border(cell)
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(size_header)
                run.font.color.rgb = _rgb("FFFFFF")

    for i in range(df.shape[0]):
        for j in range(df.shape[1]):
            cell = table.cell(i + 1, j)
            valor = df.iat[i, j]
            cell.text = "" if pd.isna(valor) else str(valor)
            _set_cell_border(cell)
            if j == 2:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(menta)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb("FFFFFF")
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
                if not p.runs:
                    p.add_run().text = ""
                for run in p.runs:
                    run.font.size = Pt(size_body)

    if w:
        total_in = sum(w)
        for j, width_in in enumerate(w):
            table.columns[j].width = Emu(int(width * (width_in / total_in)))

    return table


# ---------------------------------------------------------------------------
# Tarjetas institucionales (equivalente a rvg::dml + grid, con formas nativas)
# ---------------------------------------------------------------------------

def fmt_delta(x) -> dict:
    if x is None or (isinstance(x, float) and np.isnan(x)):
        return {"label": "s/d", "col": COL_MUTED, "icon": ""}
    if x > 0:
        return {"label": f"+{x:.0f}%", "col": COL_VERDE, "icon": "▲ "}
    if x < 0:
        return {"label": f"{x:.0f}%", "col": COL_GUINDA, "icon": "▼ "}
    return {"label": "0%", "col": COL_MUTED, "icon": "• "}


def elige_acento(
    var_2025, var_2024,
    verde: str = COL_VERDE_CHILLON, amarillo: str = COL_AMARILLO_CHILLON, rojo: str = COL_ROJO_CHILLON,
) -> str:
    v25_neg = pd.notna(var_2025) and var_2025 < 0
    v24_neg = pd.notna(var_2024) and var_2024 < 0
    if v25_neg and v24_neg:
        return rojo
    if v25_neg != v24_neg:
        return amarillo
    return verde


def _card_textbox(slide, left, top, width, height, x_frac, y_top_frac, h_frac, parts, size, align=PP_ALIGN.LEFT):
    box_left = left + int(width * x_frac)
    box_top = top + int(height * y_top_frac)
    box_width = width - int(width * x_frac) - int(width * 0.04)
    box_height = int(height * h_frac)
    tb = slide.shapes.add_textbox(box_left, box_top, box_width, box_height)
    tf = tb.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    for text, color, bold in parts:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = "Calibri"
        run.font.color.rgb = _rgb(color)
    return tb


def dibujar_card(
    slide, left, top, width, height,
    numero, titulo, var_vs_2025, var_vs_2024,
    acento: str = COL_VERDE, size_num: float = 30, size_titulo: float = 13, size_delta: float = 10.5,
    lugar_partes: list | None = None, size_lugar: float = 9,
):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = _rgb("FFFFFF")
    card.line.color.rgb = _rgb(COL_BORDE)
    card.line.width = Pt(1)
    card.shadow.inherit = False
    card.text_frame.paragraphs[0].text = ""
    card.adjustments[0] = 0.08

    barra = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left + int(width * 0.02), top + int(height * 0.05),
        int(width * 0.02), int(height * 0.90),
    )
    barra.fill.solid()
    barra.fill.fore_color.rgb = _rgb(acento)
    barra.line.fill.background()
    barra.shadow.inherit = False

    _card_textbox(slide, left, top, width, height, 0.08, 0.02, 0.28, [(fmt_num(numero), COL_DORADO, True)], size_num)
    _card_textbox(slide, left, top, width, height, 0.08, 0.30, 0.16, [(str(titulo), COL_TEXTO, True)], size_titulo)

    d25 = fmt_delta(var_vs_2025)
    d24 = fmt_delta(var_vs_2024)

    _card_textbox(
        slide, left, top, width, height, 0.08, 0.47, 0.14,
        [(f"{d25['icon']}vs 2025 ", COL_MUTED, False), (d25["label"], d25["col"], True)],
        size_delta,
    )
    _card_textbox(
        slide, left, top, width, height, 0.08, 0.61, 0.14,
        [(f"{d24['icon']}vs 2024 ", COL_MUTED, False), (d24["label"], d24["col"], True)],
        size_delta,
    )
    if lugar_partes:
        _card_textbox(
            slide, left, top, width, height, 0.08, 0.78, 0.18,
            lugar_partes,
            size_lugar,
        )


def _agregar_valuebox(slide, layout, ph_name, df, metrica, titulo, lugar_partes=None):
    ph = find_placeholder(layout, slide, ph_name)
    if ph is None:
        return
    left, top, width, height = _bbox(ph)
    _remove(ph)

    numero = valor_anio_col(df, metrica, 2026)
    numero = 0 if pd.isna(numero) else numero

    var_2025 = valor_anio_col(df, f"var_2026_vs_2025_{metrica}", 2026)
    var_2024 = valor_anio_col(df, f"var_2026_vs_2024_{metrica}", 2026)

    acento = elige_acento(var_2025, var_2024)

    dibujar_card(
        slide, left, top, width, height,
        numero=numero, titulo=titulo,
        var_vs_2025=var_2025, var_vs_2024=var_2024,
        acento=acento, lugar_partes=lugar_partes,
    )


def agregar_valueboxes(slide, layout, metricas_vb, datos_consulta_funcion, datos_curps, rankings=None):
    rankings = rankings or {}
    for i, metrica in enumerate(metricas_vb, start=1):
        lugares = rankings.get("lugares", {})
        lugar_partes = None
        if metrica in lugares:
            lugar_partes = [
                (f"{lugares[metrica]}/{rankings['total']}", COL_DORADO, True),
                ("  en su categoría", COL_MUTED, False),
            ]
        _agregar_valuebox(slide, layout, f"arriba {i}", datos_consulta_funcion, metrica,
                          MAPA_TITULOS_CONSULTAS[metrica], lugar_partes=lugar_partes)
        _agregar_valuebox(slide, layout, f"abajo {i}", datos_curps, metrica, MAPA_TITULOS_CURP[metrica])


# ---------------------------------------------------------------------------
# Ranking por categoría gerencial
# ---------------------------------------------------------------------------
CATEGORIAS_EXCLUIR_RANKING = {
    "almacenes", "laboratorios", "oficinas administrativas",
    "otros establecimientos de apoyo", "revisar",
}

MAPA_METRICA_CUBOS = {
    "consulta_gral": "consulta_general",
    "consulta_esp": "consulta_especialidad",
    "qx": "procedimientos_qx",
    "egresos": "egresos",
}

MAPA_META_COL = {
    "consulta_gral": "meta_general_anual",
    "consulta_esp": "meta_especialidad_anual",
    "qx": "meta_cirugia_anual",
    "egresos": "meta_egresos_anual",
}


def calcular_rankings_categoria(codigo_clues: str, clues_info: pd.DataFrame,
                                categoria_gerencial_df: pd.DataFrame, cubos_parquet_path: str,
                                metas: pd.DataFrame | None = None, anio: int = 2026) -> dict | None:
    """Ranking del CLUES dentro de su categoria_gerencial_nueva, por cada tipo
    de productividad, sobre el acumulado del año en curso. Devuelve None si
    no aplica (agregados como NACIONAL/entidad, o CLUES sin categoria).

    Si se pasa `metas`, solo se rankean los tipos de productividad para los
    que esta CLUES tiene meta asignada (ej. una unidad móvil que solo tiene
    meta de consulta general no debe salir rankeada en qx/egresos aunque
    tenga algo de actividad ahí)."""
    clues_validas = set(clues_info["clues_imb"].dropna().unique())
    if codigo_clues not in clues_validas:
        return None

    metricas_con_meta = None
    if metas is not None:
        metas_filtrado = metas[metas["clues_imb"] == codigo_clues]
        metricas_con_meta = {
            metrica for metrica, col_meta in MAPA_META_COL.items()
            if col_meta in metas_filtrado.columns and metas_filtrado[col_meta].fillna(0).sum() > 0
        }

    cat = categoria_gerencial_df.dropna(subset=["categoria_gerencial_nueva"]).copy()
    cat = cat[~cat["categoria_gerencial_nueva"].str.strip().str.lower().isin(CATEGORIAS_EXCLUIR_RANKING)]

    fila_categoria = cat.loc[cat["clues_imb"] == codigo_clues, "categoria_gerencial_nueva"]
    if fila_categoria.empty:
        return None
    categoria = fila_categoria.iloc[0]

    cubos = pd.read_parquet(cubos_parquet_path)
    cubos = cubos[cubos["clues"].isin(clues_validas) & (cubos["anio"].astype(str) == str(anio))]
    if cubos.empty:
        return None

    totales = cubos.groupby("clues", as_index=False).agg(
        consulta_general=("consultas_generales", "sum"),
        consulta_especialidad=("consultas_de_especialidad", "sum"),
        procedimientos_qx=("procedimientos_quirurgicos", "sum"),
        egresos=("egresos", "sum"),
    )

    pares_categoria = totales.merge(cat, left_on="clues", right_on="clues_imb", how="inner")
    pares_categoria = pares_categoria[pares_categoria["categoria_gerencial_nueva"] == categoria]

    total_categoria = len(pares_categoria)
    if total_categoria == 0 or codigo_clues not in pares_categoria["clues"].values:
        return None

    lugares = {}
    for metrica_card, col_cubos in MAPA_METRICA_CUBOS.items():
        if metricas_con_meta is not None and metrica_card not in metricas_con_meta:
            continue
        rangos = pares_categoria[col_cubos].rank(method="min", ascending=False)
        fila = rangos[pares_categoria["clues"] == codigo_clues]
        if not fila.empty:
            lugares[metrica_card] = int(fila.iloc[0])

    if not lugares:
        return None

    return {"categoria": categoria, "total": total_categoria, "lugares": lugares}


# ---------------------------------------------------------------------------
# Helpers de datos (equivalentes 1:1 a utils_crear_pptx.R)
# ---------------------------------------------------------------------------

def obtener_col(df: pd.DataFrame, nombre: str, default=0) -> pd.Series:
    if nombre in df.columns:
        return df[nombre]
    return pd.Series([default] * len(df), index=df.index)


def rellenar_anios(df: pd.DataFrame, anios) -> pd.DataFrame:
    base = pd.DataFrame({"anio": list(anios)})
    out = base.merge(df, on="anio", how="left")
    numeric_cols = out.select_dtypes(include="number").columns
    out[numeric_cols] = out[numeric_cols].fillna(0)
    return out


def asegurar_columnas(df: pd.DataFrame, cols) -> pd.DataFrame:
    df = df.copy()
    for c in cols:
        if c not in df.columns:
            df[c] = 0
    return df


def valor_anio_col(df: pd.DataFrame, columna: str, anio_objetivo):
    if columna not in df.columns:
        return np.nan
    val = df.loc[df["anio"] == anio_objetivo, columna]
    if val.empty:
        return np.nan
    return val.iloc[0]


def hay_indicador_2026(df: pd.DataFrame, columna: str) -> bool:
    val = valor_anio_col(df, columna, 2026)
    return bool(pd.notna(val) and val != 0)


def tiene_dato_2026(df: pd.DataFrame, columna: str) -> bool:
    if columna not in df.columns:
        return False
    valores = df.loc[df["anio"] == 2026, columna]
    if valores.empty:
        return False
    return bool(((valores.notna()) & (valores != 0)).any())


def definir_layout_valueboxes(datos_consulta_funcion: pd.DataFrame) -> dict:
    hay_general = tiene_dato_2026(datos_consulta_funcion, "consulta_gral")
    hay_esp = tiene_dato_2026(datos_consulta_funcion, "consulta_esp")
    hay_qx = tiene_dato_2026(datos_consulta_funcion, "qx")
    hay_egresos = tiene_dato_2026(datos_consulta_funcion, "egresos")

    if hay_general and not hay_esp and not hay_qx and not hay_egresos:
        return {"layout": "2_valueboxes", "metricas": ["consulta_gral"]}
    if hay_general and hay_esp and not hay_qx and not hay_egresos:
        return {"layout": "6_valueboxes", "metricas": ["total_consultas", "consulta_gral", "consulta_esp"]}
    if hay_general and hay_esp and hay_qx and not hay_egresos:
        return {"layout": "8_valueboxes", "metricas": ["total_consultas", "consulta_gral", "consulta_esp", "qx"]}
    if hay_general and hay_esp and hay_qx and hay_egresos:
        return {
            "layout": "10_valueboxes",
            "metricas": ["total_consultas", "consulta_gral", "consulta_esp", "qx", "egresos"],
        }

    metricas_presentes = []
    if hay_general:
        metricas_presentes.append("consulta_gral")
    if hay_esp:
        metricas_presentes.append("consulta_esp")
    if hay_qx:
        metricas_presentes.append("qx")
    if hay_egresos:
        metricas_presentes.append("egresos")
    if len(metricas_presentes) >= 2:
        metricas_presentes = ["total_consultas"] + metricas_presentes

    n = len(metricas_presentes)
    if n <= 1:
        layout_fallback = "2_valueboxes"
    elif n == 3:
        layout_fallback = "6_valueboxes"
    elif n == 4:
        layout_fallback = "8_valueboxes"
    else:
        layout_fallback = "10_valueboxes"

    return {"layout": layout_fallback, "metricas": metricas_presentes}


def armar_tabla_dinamica(df: pd.DataFrame, indicadores: list, etiquetas: list, mes_nombre: str) -> pd.DataFrame:
    filas = []
    col_25 = f"{mes_nombre} 2025"
    col_26 = f"{mes_nombre} 2026"
    for ind, etq in zip(indicadores, etiquetas):
        v25 = valor_anio_col(df, ind, 2025)
        v26 = valor_anio_col(df, ind, 2026)

        if pd.isna(v25) or v25 == 0:
            crecimiento_str = "s/d"
        else:
            crecimiento = round((v26 - v25) / v25 * 100)
            crecimiento_str = f"+{crecimiento} %" if crecimiento > 0 else f"{crecimiento} %"

        v25_str = "s/d" if pd.isna(v25) else f"{v25:,.0f}"
        v26_str = "s/d" if pd.isna(v26) else f"{v26:,.0f}"

        filas.append({"indicador": etq, col_25: v25_str, col_26: v26_str, "Crecimiento anual": crecimiento_str})

    return pd.DataFrame(filas, columns=["indicador", col_25, col_26, "Crecimiento anual"])


# ---------------------------------------------------------------------------
# Función principal
# ---------------------------------------------------------------------------

def crear_reporte_productividad(
    codigo_clues: str,
    clues_info: pd.DataFrame,
    metas: pd.DataFrame,
    historicos: pd.DataFrame,
    procedimientos_personas: pd.DataFrame,
    ruta_master: str,
) -> Presentation:
    fecha_corte = calcular_fecha_corte()
    fecha_portada = fecha_larga_es(fecha_corte)
    fecha_corte_15 = pd.Timestamp(year=fecha_corte.year, month=fecha_corte.month, day=15)
    fecha_fin_graf = fecha_corte_15

    historicos = historicos.copy()
    historicos["fecha"] = pd.to_datetime(historicos["fecha"])
    historicos["anio"] = historicos["fecha"].dt.year

    procedimientos_personas = procedimientos_personas.copy()
    procedimientos_personas["fecha"] = procedimientos_personas["fecha"].astype(float).astype(int)
    procedimientos_personas["tipo_procedimiento"] = (
        procedimientos_personas["tipo_procedimiento"].map(_TIPO_PROC_PPTX_MAP).fillna(procedimientos_personas["tipo_procedimiento"])
    )

    hay_fila_2026 = bool((procedimientos_personas["fecha"] == 2026).any())

    clues_info_filtrado = clues_info[clues_info["clues_imb"] == codigo_clues]
    if clues_info_filtrado.empty:
        raise ValueError("No se encontró el codigo_clues en clues_info.")

    metas_filtrado = metas[metas["clues_imb"] == codigo_clues]
    if metas_filtrado.empty:
        metas_filtrado = pd.DataFrame([{
            "clues_imb": codigo_clues,
            "meta_general_anual": 0, "meta_especialidad_anual": 0,
            "meta_cirugia_anual": 0, "meta_egresos_anual": 0,
        }])

    meta_total_consultas = (
        metas_filtrado["meta_general_anual"].fillna(0).sum() + metas_filtrado["meta_especialidad_anual"].fillna(0).sum()
    )
    meta_qx = metas_filtrado["meta_cirugia_anual"].fillna(0).sum()

    prs = Presentation(ruta_master)
    master = prs.slide_masters[0]

    # Ranking por categoría gerencial (solo se usa en las tarjetas de valor,
    # no en la portada) -------------------------------------------------
    categoria_gerencial_df = load_categoria_gerencial_nueva()
    ranking = calcular_rankings_categoria(
        codigo_clues, clues_info, categoria_gerencial_df, str(CUBOS_PARQUET), metas=metas_filtrado,
    )

    # Portada ------------------------------------------------------------
    slide, layout = add_slide(prs, master, "Portada 3")
    nombre_unidad = str(clues_info_filtrado["nombre_de_la_unidad"].iloc[0]).title()
    set_placeholder_text(
        layout, slide, "Título 1",
        f"Reporte de productividad médica\n{nombre_unidad} ({clues_info_filtrado['clues_imb'].iloc[0]})",
    )
    set_placeholder_text(layout, slide, "Marcador de contenido 2", fecha_portada)

    # Datos base -----------------------------------------------------------
    cols_metricas = ["consulta_gral", "consulta_esp", "qx", "total_consultas", "egresos"]

    historicos_agg = asegurar_columnas(
        historicos, ["consulta_general", "consulta_especialidad", "procedimientos_qx", "consulta_total", "egresos"]
    )
    datos_anual = historicos_agg.groupby("anio", as_index=False).agg(
        consulta_gral_anual=("consulta_general", "sum"),
        consulta_esp_anual=("consulta_especialidad", "sum"),
        qx_anual=("procedimientos_qx", "sum"),
        total_consultas_anual=("consulta_total", "sum"),
        egresos_anual=("egresos", "sum"),
    )
    datos_anual = rellenar_anios(datos_anual, range(2020, 2027))

    datos_consulta_funcion = (
        procedimientos_personas.rename(columns={"fecha": "anio"})
        .pivot_table(index="anio", columns="tipo_procedimiento", values="procedimientos", aggfunc="sum")
        .reset_index()
    )
    datos_consulta_funcion = rellenar_anios(datos_consulta_funcion, range(2024, 2027))
    datos_consulta_funcion = asegurar_columnas(datos_consulta_funcion, cols_metricas)
    datos_consulta_funcion = datos_consulta_funcion.merge(datos_anual, on="anio", how="left")
    numeric_cols = datos_consulta_funcion.select_dtypes(include="number").columns
    datos_consulta_funcion[numeric_cols] = datos_consulta_funcion[numeric_cols].fillna(0)

    datos_consulta_funcion["total_consultas_meta"] = np.where(
        datos_consulta_funcion["anio"] == 2026, meta_total_consultas, datos_consulta_funcion["total_consultas_anual"],
    )
    datos_consulta_funcion["qx_meta"] = np.where(
        datos_consulta_funcion["anio"] == 2026, meta_qx, datos_consulta_funcion["qx_anual"],
    )

    datos_curps = (
        procedimientos_personas.rename(columns={"fecha": "anio"})
        .pivot_table(index="anio", columns="tipo_procedimiento", values="personas", aggfunc="sum")
        .reset_index()
    )
    datos_curps = rellenar_anios(datos_curps, range(2024, 2027))
    datos_curps = asegurar_columnas(datos_curps, cols_metricas)
    numeric_cols = datos_curps.select_dtypes(include="number").columns
    datos_curps[numeric_cols] = datos_curps[numeric_cols].fillna(0)

    for col in cols_metricas:
        for ref in (2024, 2025):
            nombre_var = f"var_2026_vs_{ref}_{col}"
            for df_ in (datos_consulta_funcion, datos_curps):
                valor_2026 = valor_anio_col(df_, col, 2026)
                valor_ref = valor_anio_col(df_, col, ref)
                # R permite Inf/-Inf en este cálculo sin truronar (valor_2026==0 con
                # valor_ref!=0); en Python hay que evitarlo explícitamente para no
                # crashear con OverflowError al redondear un infinito.
                if pd.notna(valor_2026) and pd.notna(valor_ref) and valor_ref != 0 and valor_2026 != 0:
                    computed = round(100 * (1 - (valor_ref / valor_2026)))
                else:
                    computed = 0
                df_[nombre_var] = np.where(df_["anio"] == 2026, computed, 0)

    # Value boxes ------------------------------------------------------------
    config_vb = definir_layout_valueboxes(datos_consulta_funcion)
    layout_vb, metricas_vb = config_vb["layout"], config_vb["metricas"]

    slide, layout = add_slide(prs, master, layout_vb)
    set_placeholder_text(layout, slide, "Título 1", "Productividad IMSS Bienestar")
    set_placeholder_text(layout, slide, "fecha", f"Del 01 de enero al {fecha_portada}")
    agregar_valueboxes(slide, layout, metricas_vb, datos_consulta_funcion, datos_curps, rankings=ranking)

    # Diapo 3: 2024-2026, solo si hay 2026 ------------------------------------
    if hay_fila_2026:
        d = datos_consulta_funcion.copy()
        d["anio_num"] = d["anio"].astype(int)
        for c in ["consulta_gral", "consulta_esp", "qx", "egresos", "consulta_gral_anual", "consulta_esp_anual", "qx_anual"]:
            d[c] = d[c].fillna(0)

        d["total_consultas"] = np.select(
            [(d["consulta_gral"] > 0) & (d["consulta_esp"] > 0), d["consulta_gral"] > 0, d["consulta_esp"] > 0],
            [d["consulta_gral"] + d["consulta_esp"], d["consulta_gral"], d["consulta_esp"]],
            default=0,
        )

        d["total_consultas_meta"] = np.select(
            [
                (d["anio_num"] == 2026) & d["total_consultas_meta"].notna(),
                d["anio_num"] == 2026,
                (d["consulta_gral_anual"] > 0) & (d["consulta_esp_anual"] > 0),
                d["consulta_gral_anual"] > 0,
                d["consulta_esp_anual"] > 0,
                d["total_consultas_anual"] > 0,
            ],
            [
                d["total_consultas_meta"], d["total_consultas"],
                d["consulta_gral_anual"] + d["consulta_esp_anual"], d["consulta_gral_anual"],
                d["consulta_esp_anual"], d["total_consultas_anual"],
            ],
            default=d["total_consultas"],
        )

        d["qx_meta"] = np.select(
            [(d["anio_num"] == 2026) & d["qx_meta"].notna(), d["anio_num"] == 2026, d["qx_anual"] > 0],
            [d["qx_meta"], d["qx"], d["qx_anual"]],
            default=d["qx"],
        )

        datos_2024_2026 = d[d["anio_num"].isin([2024, 2025, 2026])]

        hay_consultas_2024_2026 = bool((datos_2024_2026["total_consultas"] > 0).any())
        hay_qx_2024_2026 = bool(((datos_2024_2026["qx"] > 0) | (datos_2024_2026["egresos"] > 0)).any())

        if hay_consultas_2024_2026:
            indicadores_consulta, etiquetas_consulta = [], []
            if hay_indicador_2026(datos_2024_2026, "consulta_gral"):
                indicadores_consulta.append("consulta_gral")
                etiquetas_consulta.append("Consultas generales")
            if hay_indicador_2026(datos_2024_2026, "consulta_esp"):
                indicadores_consulta.append("consulta_esp")
                etiquetas_consulta.append("Consultas de especialidad*")

            tabla_consultas = armar_tabla_dinamica(datos_2024_2026, indicadores_consulta, etiquetas_consulta, "Acumulado")

            if hay_qx_2024_2026:
                indicadores_proc, etiquetas_proc = [], []
                if hay_indicador_2026(datos_2024_2026, "qx"):
                    indicadores_proc.append("qx")
                    etiquetas_proc.append("Procedimientos quirúrgicos")
                if hay_indicador_2026(datos_2024_2026, "egresos"):
                    indicadores_proc.append("egresos")
                    etiquetas_proc.append("Egresos")

                tabla_proc = armar_tabla_dinamica(datos_2024_2026, indicadores_proc, etiquetas_proc, "Acumulado")

                slide, layout = add_slide(prs, master, "Historico consultas y procedimientos")
                set_placeholder_text(layout, slide, "Título 1", "Productividad IMSS Bienestar")
                place_native_chart(slide, layout, "Grafica 1", dibujar_grafica_planeacion_2024_2026,
                                   datos_2024_2026, "total_consultas_meta", "total_consultas", "Consultas totales")
                place_native_chart(slide, layout, "Grafica 2", dibujar_grafica_planeacion_2024_2026,
                                   datos_2024_2026, "qx_meta", "qx", "Procedimientos quirúrgicos")
                place_table(slide, layout, "tabla_1", tabla_consultas, w=[2.70, 0.90, 0.90, 0.80], size_header=8, size_body=8)
                place_table(slide, layout, "tabla_2", tabla_proc, w=[2.70, 0.90, 0.90, 0.80], size_header=8, size_body=8)
                set_placeholder_text(layout, slide, "fecha", f"Del 01 de enero al {fecha_portada}")
            else:
                slide, layout = add_slide(prs, master, "Historico consultas")
                set_placeholder_text(layout, slide, "Título 1", "Productividad IMSS Bienestar")
                place_native_chart(slide, layout, "Grafica 1", dibujar_grafica_planeacion_2024_2026,
                                   datos_2024_2026, "total_consultas_meta", "total_consultas", "Consultas totales")
                place_table(slide, layout, "tabla_1", tabla_consultas, w=[4.60, 1.35, 1.35, 1.40], size_header=11, size_body=10)
                set_placeholder_text(layout, slide, "fecha", f"Del 01 de enero al {fecha_portada}")

    # Diapo 4: históricos 2020-2025 -------------------------------------------
    dh = historicos.copy()
    dh["consulta_general_tmp"] = obtener_col(dh, "consulta_general") + obtener_col(dh, "consulta_gral")
    dh["consulta_esp_tmp"] = obtener_col(dh, "consulta_especialidad") + obtener_col(dh, "consulta_esp")
    dh["procedimientos_qx_tmp"] = obtener_col(dh, "procedimientos_qx")
    dh["egresos_tmp"] = obtener_col(dh, "egresos")

    mask = (
        (dh["fecha"].dt.year >= 2020)
        & (dh["fecha"].dt.year <= 2025)
        & (dh["fecha"].dt.strftime("%m-%d") <= fecha_corte_15.strftime("%m-%d"))
    )
    dh = dh[mask].copy()
    dh["anio"] = dh["fecha"].dt.year.astype(str)

    resumen = dh.groupby("anio").agg(
        consulta_gral=("consulta_general_tmp", "sum"),
        consulta_esp=("consulta_esp_tmp", "sum"),
        qx=("procedimientos_qx_tmp", "sum"),
        egresos=("egresos_tmp", "sum"),
    ).reset_index()

    resumen["total_consultas"] = np.select(
        [(resumen["consulta_gral"] > 0) & (resumen["consulta_esp"] > 0), resumen["consulta_gral"] > 0, resumen["consulta_esp"] > 0],
        [resumen["consulta_gral"] + resumen["consulta_esp"], resumen["consulta_gral"], resumen["consulta_esp"]],
        default=0,
    )

    datos_anual_str = datos_anual.copy()
    datos_anual_str["anio"] = datos_anual_str["anio"].astype(str)
    resumen = resumen.merge(datos_anual_str, on="anio", how="left")
    for c in ["consulta_gral_anual", "consulta_esp_anual", "qx_anual"]:
        resumen[c] = resumen[c].fillna(0)

    resumen["total_consultas_anual"] = np.select(
        [
            (resumen["consulta_gral_anual"] > 0) & (resumen["consulta_esp_anual"] > 0),
            resumen["consulta_gral_anual"] > 0,
            resumen["consulta_esp_anual"] > 0,
        ],
        [resumen["consulta_gral_anual"] + resumen["consulta_esp_anual"], resumen["consulta_gral_anual"], resumen["consulta_esp_anual"]],
        default=0,
    )
    datos_historicos_2020_2025 = resumen.sort_values("anio")

    hay_consultas_2020 = bool((datos_historicos_2020_2025["total_consultas"] > 0).any())
    hay_qx_2020 = bool(
        ((datos_historicos_2020_2025["qx"] > 0) | (datos_historicos_2020_2025["egresos"] > 0)).any()
    )

    if hay_consultas_2020:
        if hay_qx_2020:
            slide, layout = add_slide(prs, master, "1_Historico consultas y procedimientos")
            set_placeholder_text(layout, slide, "Título 1", "Productividad IMSS Bienestar")
            place_native_chart(slide, layout, "Grafica 1", dibujar_grafica_planeacion_historica,
                               datos_historicos_2020_2025, "total_consultas_anual", "total_consultas",
                               "Consultas totales")
            place_native_chart(slide, layout, "Grafica 2", dibujar_grafica_planeacion_historica,
                               datos_historicos_2020_2025, "qx_anual", "qx", "Procedimientos quirúrgicos")
            set_placeholder_text(layout, slide, "fecha", f"Del 01 de enero al {fecha_portada}")
        else:
            slide, layout = add_slide(prs, master, "1_Historico consultas")
            set_placeholder_text(layout, slide, "Título 1", "Productividad IMSS Bienestar")
            place_native_chart(slide, layout, "Grafica 1", dibujar_grafica_planeacion_historica,
                               datos_historicos_2020_2025, "total_consultas_anual", "total_consultas",
                               "Consultas totales")
            set_placeholder_text(layout, slide, "fecha", f"Del 01 de enero al {fecha_portada}")

    # Diapo 5 ------------------------------------------------------------
    serie_mensual_consultas = historicos.copy()
    serie_mensual_consultas["fecha"] = serie_mensual_consultas["fecha"].values.astype("datetime64[M]")
    serie_mensual_consultas = serie_mensual_consultas.dropna(subset=["fecha", "consulta_total"])
    serie_mensual_consultas = (
        serie_mensual_consultas.groupby("fecha", as_index=False)["consulta_total"]
        .sum()
        .rename(columns={"consulta_total": "consultas_totales"})
        .sort_values("fecha")
    )

    if len(serie_mensual_consultas) > 0:
        slide, layout = add_slide(prs, master, "Una grafica")
        set_placeholder_text(layout, slide, "Título 1", "Consultas totales por mes (2022-2026)")
        place_native_chart(
            slide, layout, "ft", dibujar_grafica_consultas_periodos,
            serie_mensual_consultas, fecha_inicio="2022-08-01", fecha_fin=str(fecha_fin_graf.date()),
            titulo=f"Consultas totales del IMSS Bienestar (agosto 2022 – {mes_anio_es(fecha_fin_graf)})",
        )

    # Diapo 6 ------------------------------------------------------------
    if hay_fila_2026 and hay_indicador_2026(datos_consulta_funcion, "qx"):
        serie_mensual_pq = historicos.copy()
        serie_mensual_pq["fecha"] = serie_mensual_pq["fecha"].values.astype("datetime64[M]")
        serie_mensual_pq = serie_mensual_pq.dropna(subset=["fecha", "procedimientos_qx"])
        serie_mensual_pq = (
            serie_mensual_pq.groupby("fecha", as_index=False)["procedimientos_qx"]
            .sum()
            .rename(columns={"procedimientos_qx": "consultas_totales"})
            .sort_values("fecha")
        )

        if len(serie_mensual_pq) > 0:
            slide, layout = add_slide(prs, master, "Una grafica")
            set_placeholder_text(layout, slide, "Título 1", "Procedimientos quirúrgicos por mes (2022-2026)")
            place_native_chart(
                slide, layout, "ft", dibujar_grafica_consultas_periodos,
                serie_mensual_pq, fecha_inicio="2022-08-01", fecha_fin=str(fecha_fin_graf.date()),
                titulo=f"Procedimientos quirúrgicos del IMSS Bienestar (agosto 2022 – {mes_anio_es(fecha_fin_graf)})",
            )

    return prs
